"""File processing and discovery for peekdocs."""

# Filename prefixes for the three kinds of result reports peekdocs writes.
# Used everywhere that needs to discover, list, or clean up result files:
# the CLI (--list-files / --clear / --clear-all), the GUI's "View All
# peekdocs Files" / "Delete on Close" / per-search stale-file cleanup, and
# the scanner's own skip-list (so a search doesn't recurse into its own
# prior reports).
RESULT_FILE_PREFIXES = (
    "peekdocs_standard_results",
    "peekdocs_regex_results",
    "peekdocs_suite_results",
)


def is_peekdocs_internal_file(basename):
    """Return True if *basename* is a peekdocs-generated file.

    The naming convention documented in README, USER_GUIDE, and
    SECURITY.md reserves the `peekdocs_` prefix (visible files) and
    `.peekdocs` prefix (hidden user-state and per-folder dotfiles)
    for peekdocs's own outputs. All Search workflows
    (Standard / Suites / Regex) skip such files at discovery time
    so peekdocs never finds its own reports, state files, or tools-
    menu outputs as user content. Used by both the scanner's
    discover_files loop and the GUI's Find Duplicates tool so the
    two surfaces agree on what counts as a peekdocs file.
    """
    return basename.startswith(("peekdocs_", ".peekdocs"))

import csv
import glob
import os
import re
import stat
import warnings
from copy import deepcopy
from html.parser import HTMLParser
from itertools import product

# Suppress noisy openpyxl warnings about malformed dates, unsupported
# extensions, etc. These don't affect search results — they're cosmetic.
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")

import fitz
from docx import Document
from odf.opendocument import load as load_odt
from odf.text import P as OdtParagraph
from odf.table import Table as OdfTable, TableRow as OdfTableRow, TableCell as OdfTableCell
from odf import teletype
from openpyxl import load_workbook
from striprtf.striprtf import rtf_to_text
from pptx import Presentation as PptxPresentation
import ebooklib
from ebooklib import epub

import email as email_mod
import tempfile
import zipfile
import tarfile
import olefile
import xlrd
import extract_msg
try:
    import pypff
except ImportError:
    pypff = None
import py7zr
import rarfile

from peekdocs.constants import SUPPORTED_TYPES, OCR_IMAGE_TYPES, FUZZY_THRESHOLD
from peekdocs.range_query import line_matches_content_ranges, file_matches_metadata_ranges, file_matches_filename_ranges


def apply_context(all_lines, match_indices, before, after):
    """Expand match indices with before/after context and return merged groups."""
    if not match_indices:
        return []

    total = len(all_lines)
    ranges = []
    for idx in sorted(match_indices):
        start = max(0, idx - before)
        end = min(total - 1, idx + after)
        ranges.append((start, end))

    merged = []
    for start, end in ranges:
        if merged and start <= merged[-1][1] + 1:
            merged[-1] = (merged[-1][0], max(merged[-1][1], end))
        else:
            merged.append([start, end])

    groups = []
    for start, end in merged:
        group = []
        for i in range(start, end + 1):
            line_num, text = all_lines[i]
            is_match = i in match_indices
            group.append((line_num, text, is_match))
        groups.append(group)

    return groups


def _wildcard_to_regex(term):
    """Convert a wildcard pattern (* and ?) to a regex pattern.

    * matches zero or more non-whitespace characters (including punctuation).
    ? matches exactly one non-whitespace character.
    """
    escaped = re.escape(term)
    escaped = escaped.replace(r'\*', r'\S*')
    escaped = escaped.replace(r'\?', r'\S')
    return r'\b' + escaped


def _whole_word_pattern(term):
    """Build a regex pattern for whole-word matching.

    Only adds \\b where the term starts/ends with a word character.
    This prevents matching failures when the term contains leading or
    trailing punctuation (e.g., 'output_min;' ends with ';' which is
    not a word character, so \\b after it would never match).
    """
    escaped = re.escape(term)
    prefix = r'\b' if re.match(r'\w', term) else ''
    suffix = r'\b' if re.search(r'\w$', term) else ''
    return prefix + escaped + suffix


def _ocr_image(image):
    """Run OCR on a PIL Image and return extracted text."""
    import pytesseract
    from peekdocs.paths import find_tesseract

    # pytesseract shells out to the ``tesseract`` binary and by default
    # relies on the current process's PATH to find it. That breaks under
    # macOS GUI launches (Finder / Dock / Spotlight give a stripped PATH
    # that omits /opt/homebrew/bin) even when Tesseract is installed via
    # Homebrew. find_tesseract() checks well-known install locations too,
    # so we pin pytesseract to the absolute path we located.
    tesseract_path = find_tesseract()
    if tesseract_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
    return pytesseract.image_to_string(image)


def _extract_lines(filepath, use_ocr=False, ocr_func=None):
    """Extract text lines from a file.

    Returns a list of (line_num, text) tuples. The meaning of line_num
    varies by file type (line number, page number, row number, etc.).
    Raises an exception on error — the caller decides how to handle it.
    """
    if ocr_func is None:
        ocr_func = _ocr_image

    ext = os.path.splitext(filepath)[1].lower()
    # Handle special filenames with no standard extension
    basename_lower = os.path.basename(filepath).lower()
    if not ext and basename_lower in (".env", "dockerfile"):
        ext = f".{basename_lower}" if not basename_lower.startswith(".") else basename_lower
    all_lines = []

    if ext == ".docx":
        doc = Document(filepath)
        all_lines = [(i, para.text) for i, para in enumerate(doc.paragraphs, start=1)]

    elif ext == ".doc":
        # Word 97-2003 binary format — extract text from OLE compound document
        ole = olefile.OleFileIO(filepath)
        if ole.exists("WordDocument"):
            stream = ole.openstream("WordDocument")
            data = stream.read()
            stream.close()
            # Try both UTF-16-LE (Unicode docs) and Latin-1 (ASCII docs),
            # use whichever produces more readable content
            candidates = []
            for encoding in ("utf-16-le", "latin-1"):
                try:
                    decoded = data.decode(encoding, errors="ignore")
                    lines = [l.strip() for l in decoded.split("\r") if l.strip() and len(l.strip()) > 2]
                    clean = [l for l in lines if sum(1 for c in l if c.isprintable()) / max(len(l), 1) > 0.7]
                    candidates.append(clean)
                except Exception:
                    candidates.append([])
            best = max(candidates, key=lambda c: sum(len(l) for l in c))
            all_lines = [(i, line) for i, line in enumerate(best, start=1)]
        ole.close()

    elif ext == ".pdf":
        doc = fitz.open(filepath)
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            if use_ocr and len((text or "").strip()) < 10:
                import io
                from PIL import Image
                pix = page.get_pixmap(dpi=300)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = ocr_func(img)
            if not text or not text.strip():
                continue
            for line in text.split("\n"):
                all_lines.append((page_num, line))
        doc.close()

    elif ext == ".csv":
        with open(filepath, newline="", encoding="utf-8", errors="replace") as csvfile:
            reader = csv.reader(csvfile)
            all_lines = [(row_num, ", ".join(row)) for row_num, row in enumerate(reader, start=1)]

    elif ext == ".tsv":
        with open(filepath, newline="", encoding="utf-8", errors="replace") as tsvfile:
            reader = csv.reader(tsvfile, delimiter="\t")
            all_lines = [(row_num, "\t".join(row)) for row_num, row in enumerate(reader, start=1)]

    elif ext == ".epub":
        book = epub.read_epub(filepath)
        line_num = 0
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            html_content = item.get_content().decode("utf-8", errors="replace")
            class _EpubTextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text_parts = []
                def handle_data(self, data):
                    self.text_parts.append(data)
            parser = _EpubTextExtractor()
            parser.feed(html_content)
            for line in "".join(parser.text_parts).split("\n"):
                stripped = line.strip()
                if stripped:
                    line_num += 1
                    all_lines.append((line_num, stripped))

    elif ext == ".odt":
        odt_doc = load_odt(filepath)
        all_lines = [(i, teletype.extractText(para)) for i, para in enumerate(odt_doc.getElementsByType(OdtParagraph), start=1)]

    elif ext == ".odp":
        odp_doc = load_odt(filepath)
        all_lines = [(i, teletype.extractText(para)) for i, para in enumerate(odp_doc.getElementsByType(OdtParagraph), start=1)]

    elif ext == ".ods":
        ods_doc = load_odt(filepath)
        row_num = 0
        for table in ods_doc.getElementsByType(OdfTable):
            for row in table.getElementsByType(OdfTableRow):
                row_num += 1
                cells = []
                for cell in row.getElementsByType(OdfTableCell):
                    cell_text = teletype.extractText(cell)
                    if cell_text:
                        cells.append(cell_text)
                row_text = ", ".join(cells)
                all_lines.append((row_num, row_text))

    elif ext in (".txt", ".md", ".json", ".xml", ".log", ".yaml", ".yml", ".toml", ".rst", ".tex", ".ini", ".cfg", ".sql", ".ics", ".vcf"):
        with open(filepath, encoding="utf-8-sig", errors="replace") as txtfile:
            all_lines = [(line_num, line.rstrip("\n")) for line_num, line in enumerate(txtfile, start=1)]

    elif ext == ".mbox":
        # Unix mailbox archive — multiple emails in one file
        import mailbox
        mbox = mailbox.mbox(filepath)
        line_num = 0
        for msg in mbox:
            for header in ("From", "To", "Subject", "Date"):
                val = msg.get(header, "")
                if val:
                    line_num += 1
                    all_lines.append((line_num, f"{header}: {val}"))
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            body = payload.decode("utf-8", errors="replace")
                            break
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    body = payload.decode("utf-8", errors="replace")
            for line in body.split("\n"):
                if line.strip():
                    line_num += 1
                    all_lines.append((line_num, line.rstrip("\r\n")))
        mbox.close()

    elif ext == ".pages":
        # Apple Pages — zip archive containing XML documents
        try:
            with zipfile.ZipFile(filepath, "r") as zf:
                line_num = 0
                for name in zf.namelist():
                    if name.endswith(".xml") or name.endswith(".html"):
                        data = zf.read(name).decode("utf-8", errors="replace")
                        class _PagesTextExtractor(HTMLParser):
                            def __init__(self):
                                super().__init__()
                                self.text_parts = []
                            def handle_data(self, data):
                                self.text_parts.append(data)
                        parser = _PagesTextExtractor()
                        parser.feed(data)
                        for line in "".join(parser.text_parts).split("\n"):
                            stripped = line.strip()
                            if stripped:
                                line_num += 1
                                all_lines.append((line_num, stripped))
        except zipfile.BadZipFile:
            pass  # Newer Pages format may use protobuf — skip gracefully

    elif ext == ".html":
        class _HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text_parts = []
            def handle_data(self, data):
                self.text_parts.append(data)
        with open(filepath, encoding="utf-8", errors="replace") as htmlfile:
            parser = _HTMLTextExtractor()
            parser.feed(htmlfile.read())
        lines = "".join(parser.text_parts).split("\n")
        all_lines = [(line_num, line.strip()) for line_num, line in enumerate(lines, start=1)]

    elif ext == ".pptx":
        prs = PptxPresentation(filepath)
        para_num = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_num += 1
                        all_lines.append((para_num, para.text))

    elif ext == ".ppt":
        # PowerPoint 97-2003 binary format — extract text from OLE streams
        ole = olefile.OleFileIO(filepath)
        if ole.exists("PowerPoint Document"):
            stream = ole.openstream("PowerPoint Document")
            data = stream.read()
            stream.close()
            line_num = 0
            i = 0
            while i < len(data) - 8:
                rec_type = int.from_bytes(data[i+2:i+4], "little")
                rec_len = int.from_bytes(data[i+4:i+8], "little")
                if rec_type == 0x0FA8 and 0 < rec_len < 100000:
                    # TextBytesAtom — ASCII text
                    text = data[i+8:i+8+rec_len].decode("latin-1", errors="ignore").strip()
                    if text:
                        for line in text.split("\r"):
                            if line.strip():
                                line_num += 1
                                all_lines.append((line_num, line.strip()))
                    i += 8 + rec_len
                elif rec_type == 0x0FA0 and 0 < rec_len < 100000:
                    # TextCharsAtom — UTF-16-LE text
                    text = data[i+8:i+8+rec_len].decode("utf-16-le", errors="ignore").strip()
                    if text:
                        for line in text.split("\r"):
                            if line.strip():
                                line_num += 1
                                all_lines.append((line_num, line.strip()))
                    i += 8 + rec_len
                else:
                    i += 1
        ole.close()

    elif ext == ".rtf":
        with open(filepath, encoding="utf-8", errors="replace") as rtffile:
            raw = rtffile.read()
        plain = rtf_to_text(raw)
        all_lines = [(line_num, line) for line_num, line in enumerate(plain.split("\n"), start=1)]

    elif ext == ".eml":
        # Standard email format (RFC 822) — one email per file
        with open(filepath, "rb") as f:
            msg = email_mod.message_from_bytes(f.read())
        line_num = 0
        # Include key headers as searchable lines
        for hdr in ("From", "To", "Cc", "Subject", "Date"):
            val = msg.get(hdr)
            if val:
                line_num += 1
                all_lines.append((line_num, f"{hdr}: {val}"))
        # Extract body text
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                if ctype == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        text = payload.decode("utf-8", errors="replace")
                        for line in text.split("\n"):
                            line_num += 1
                            all_lines.append((line_num, line.rstrip("\r")))
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                text = payload.decode("utf-8", errors="replace")
                for line in text.split("\n"):
                    line_num += 1
                    all_lines.append((line_num, line.rstrip("\r")))

    elif ext == ".msg":
        # Microsoft Outlook email format — one email per file
        msg = extract_msg.Message(filepath)
        line_num = 0
        for hdr, val in [("From", msg.sender), ("To", msg.to), ("Cc", msg.cc),
                         ("Subject", msg.subject), ("Date", msg.date)]:
            if val:
                line_num += 1
                all_lines.append((line_num, f"{hdr}: {val}"))
        body = msg.body
        if body:
            for line in body.split("\n"):
                line_num += 1
                all_lines.append((line_num, line.rstrip("\r")))
        msg.close()

    elif ext == ".pst":
        # Outlook Personal Storage Table — mailbox database with many emails
        if pypff is None:
            raise ImportError(
                "PST support requires the libpff-python package. "
                "Install it with: pip install libpff-python "
                "(requires a C compiler — see Troubleshooting in the docs)"
            )
        pst = pypff.file()
        pst.open(filepath)
        line_num = 0

        def _walk_pst_folder(folder):
            nonlocal line_num
            for i in range(folder.number_of_sub_messages):
                msg = folder.get_sub_message(i)
                for hdr, val in [("From", msg.sender_name), ("Subject", msg.subject),
                                 ("Date", str(msg.delivery_time) if msg.delivery_time else None)]:
                    if val:
                        line_num += 1
                        all_lines.append((line_num, f"{hdr}: {val}"))
                body = msg.plain_text_body
                if body:
                    text = body.decode("utf-8", errors="replace") if isinstance(body, bytes) else body
                    for line in text.split("\n"):
                        line_num += 1
                        all_lines.append((line_num, line.rstrip("\r")))
            for i in range(folder.number_of_sub_folders):
                _walk_pst_folder(folder.get_sub_folder(i))

        root = pst.get_root_folder()
        if root:
            _walk_pst_folder(root)
        pst.close()

    elif ext == ".xlsx":
        wb = load_workbook(filepath, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            for row_num, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                row_text = ", ".join(str(cell) for cell in row if cell is not None)
                all_lines.append((row_num, row_text))
        wb.close()

    elif ext == ".xls":
        # Excel 97-2003 binary format
        book = xlrd.open_workbook(filepath)
        for sheet in book.sheets():
            for row_num in range(sheet.nrows):
                vals = [str(sheet.cell_value(row_num, col)) for col in range(sheet.ncols) if sheet.cell_value(row_num, col) != ""]
                row_text = ", ".join(vals)
                if row_text.strip():
                    all_lines.append((row_num + 1, row_text))

    elif ext in OCR_IMAGE_TYPES:
        from PIL import Image
        img = Image.open(filepath)
        text = ocr_func(img)
        if text and text.strip():
            all_lines = [(line_num, line) for line_num, line in enumerate(text.split("\n"), start=1) if line.strip()]

    elif ext in (".zip", ".tar", ".gz", ".bz2", ".tgz", ".7z", ".rar"):
        # Archive formats — extract to temp dir and process each contained file
        all_supported = SUPPORTED_TYPES | OCR_IMAGE_TYPES
        _MAX_ARCHIVE_BYTES = 500 * 1024 * 1024  # 500 MB extraction limit
        with tempfile.TemporaryDirectory() as tmpdir:
            # Extract archive contents — check for zip bombs on .zip files
            if ext == ".zip":
                with zipfile.ZipFile(filepath, "r") as zf:
                    total_uncompressed = sum(i.file_size for i in zf.infolist())
                    if total_uncompressed > _MAX_ARCHIVE_BYTES:
                        raise ValueError(
                            f"Archive expands to {total_uncompressed / (1024*1024):.0f} MB, "
                            f"exceeds 500 MB safety limit. Skipping to prevent memory issues.")
                    zf.extractall(tmpdir)
            elif ext in (".tar", ".gz", ".bz2", ".tgz"):
                try:
                    with tarfile.open(filepath, "r:*") as tf:
                        # Extract files individually to handle path-too-long errors
                        for member in tf.getmembers():
                            try:
                                try:
                                    tf.extract(member, tmpdir, filter="data")
                                except TypeError:
                                    # Python < 3.11.4 doesn't support the filter parameter
                                    tf.extract(member, tmpdir)
                            except (OSError, IOError):
                                continue  # Skip files with paths too long for Windows
                except (tarfile.ReadError, tarfile.CompressionError):
                    # Not a tar archive — try as raw compressed file (e.g., plain .gz)
                    import gzip as _gzip_mod
                    if ext == ".gz":
                        basename = os.path.basename(filepath)
                        inner_name = basename[:-3] if basename.endswith(".gz") else basename + ".txt"
                        inner_path = os.path.join(tmpdir, inner_name)
                        try:
                            with _gzip_mod.open(filepath, "rb") as gz_in:
                                with open(inner_path, "wb") as gz_out:
                                    gz_out.write(gz_in.read())
                        except Exception:
                            raise  # Re-raise so the outer handler reports it
            elif ext == ".7z":
                with py7zr.SevenZipFile(filepath, "r") as sz:
                    sz.extractall(tmpdir)
            elif ext == ".rar":
                with rarfile.RarFile(filepath, "r") as rf:
                    rf.extractall(tmpdir)

            # Walk extracted files and process each supported type
            line_num = 0
            for root, _dirs, files in os.walk(tmpdir):
                for fname in sorted(files):
                    inner_ext = os.path.splitext(fname)[1].lower()
                    if inner_ext not in all_supported:
                        continue
                    # Skip nested archives to avoid infinite recursion
                    if inner_ext in (".zip", ".tar", ".gz", ".bz2", ".tgz", ".7z", ".rar"):
                        continue
                    inner_path = os.path.join(root, fname)
                    try:
                        inner_lines = _extract_lines(inner_path, use_ocr, ocr_func)
                        for _ln, text in inner_lines:
                            if text.strip():
                                line_num += 1
                                all_lines.append((line_num, f"[{fname}] {text}"))
                    except Exception:
                        continue  # Skip files that fail inside the archive

    elif ext in (".numbers", ".key"):
        # Apple iWork files — ZIP archives. Older versions contain index.xml
        # with readable text. Newer versions use protobuf (.iwa) which we
        # can't parse, but XML metadata files may still contain some text.
        # Skip non-iWork .key files (e.g., SSL certificate keys) — they're
        # not zip archives and will fail.
        if not zipfile.is_zipfile(filepath):
            return all_lines
        import xml.etree.ElementTree as _iwork_ET
        line_num = 0
        with zipfile.ZipFile(filepath, "r") as zf:
            for name in sorted(zf.namelist()):
                if not name.endswith(".xml"):
                    continue
                try:
                    tree = _iwork_ET.parse(zf.open(name))
                    for elem in tree.iter():
                        if elem.text and elem.text.strip():
                            for line in elem.text.strip().split("\n"):
                                line = line.strip()
                                if line:
                                    line_num += 1
                                    all_lines.append((line_num, line))
                        if elem.tail and elem.tail.strip():
                            for line in elem.tail.strip().split("\n"):
                                line = line.strip()
                                if line:
                                    line_num += 1
                                    all_lines.append((line_num, line))
                except Exception:
                    continue

    elif ext == ".vsdx":
        # Visio diagram — ZIP archive containing XML pages with text in shapes
        import xml.etree.ElementTree as _vsdx_ET
        line_num = 0
        with zipfile.ZipFile(filepath, "r") as zf:
            for name in sorted(zf.namelist()):
                if not name.startswith("visio/pages/page") or not name.endswith(".xml"):
                    continue
                try:
                    tree = _vsdx_ET.parse(zf.open(name))
                    ns = {"v": "http://schemas.microsoft.com/office/visio/2012/main"}
                    for text_elem in tree.iter("{http://schemas.microsoft.com/office/visio/2012/main}Text"):
                        # Collect all text content including child elements
                        full_text = "".join(text_elem.itertext()).strip()
                        if full_text:
                            for line in full_text.split("\n"):
                                line = line.strip()
                                if line:
                                    line_num += 1
                                    all_lines.append((line_num, line))
                except Exception:
                    continue

    elif ext == ".ipynb":
        # Jupyter notebook — JSON file containing code and markdown cells.
        # Extract source lines from all cells, preserving cell order.
        import json as _json_nb
        with open(filepath, encoding="utf-8", errors="replace") as nbfile:
            try:
                nb = _json_nb.load(nbfile)
            except _json_nb.JSONDecodeError:
                nb = {}
        line_num = 0
        for cell in nb.get("cells", []):
            cell_type = cell.get("cell_type", "")
            source = cell.get("source", [])
            if cell_type in ("code", "markdown", "raw"):
                for src_line in source:
                    for part in src_line.split("\n"):
                        line_num += 1
                        all_lines.append((line_num, part.rstrip("\n")))

    elif ext in SUPPORTED_TYPES:
        # Plain text fallback — source code, engineering files, and any other
        # text-based format in SUPPORTED_TYPES without a specialized parser.
        with open(filepath, encoding="utf-8-sig", errors="replace") as txtfile:
            all_lines = [(line_num, line.rstrip("\n")) for line_num, line in enumerate(txtfile, start=1)]

    return all_lines


def _search_file_lines(all_lines, file_dir, filename, config):
    """Search extracted lines for matches.

    Args:
        all_lines: list of (line_num, text) tuples from _extract_lines
        file_dir: directory containing the file
        filename: base name of the file
        config: search configuration dict

    Returns:
        (matches, skipped) where:
        - matches = list of (file_dir, filename, line_num, text) tuples
        - skipped = list of (filename, error_msg) tuples (usually empty)
    """
    search_terms = config["search_terms"]
    content_ranges = config.get("content_ranges", [])
    range_only = not search_terms and not config.get("expression_ast") and content_ranges
    use_regex = config["use_regex"]
    line_prox = config.get("line_proximity", 0)
    # Line proximity is inherently AND across lines — force OR mode for
    # per-line matching so each term is found individually, then the
    # line proximity filter enforces the cross-line AND constraint.
    match_all = config["match_all"] if line_prox == 0 else False
    use_proximity = config["use_proximity"]
    proximity = config["proximity"]
    use_context = config["use_context"]
    context_before = config["context_before"]
    context_after = config["context_after"]
    use_fuzzy = config.get("use_fuzzy", False)
    exclude_terms = config.get("exclude_terms", [])
    use_wildcard = config.get("use_wildcard", False)
    use_whole_word = config.get("use_whole_word", False)
    expression_ast = config.get("expression_ast")
    if use_wildcard and expression_ast is None:
        search_terms = [_wildcard_to_regex(t) for t in search_terms]
        if exclude_terms:
            exclude_terms = [_wildcard_to_regex(t) for t in exclude_terms]
        use_regex = True
    if use_fuzzy:
        from rapidfuzz import fuzz

    def _fuzzy_word_match(text, term):
        """Return the first word in text that fuzzy-matches term, or None."""
        for word in re.findall(r'\S+', text):
            clean = re.sub(r'^[^\w]+|[^\w]+$', '', word)
            if not clean:
                continue
            if fuzz.ratio(term.lower(), clean.lower()) >= FUZZY_THRESHOLD:
                return word
        return None

    def _single_term_matches(term, text):
        """Return True if a single term matches the text, respecting the current mode."""
        if use_wildcard:
            pattern = _wildcard_to_regex(term)
            return bool(re.search(pattern, text, re.IGNORECASE))
        elif use_regex:
            return bool(re.search(term, text, re.IGNORECASE))
        elif use_fuzzy:
            return _fuzzy_word_match(text, term) is not None
        elif use_whole_word:
            return bool(re.search(_whole_word_pattern(term), text, re.IGNORECASE))
        else:
            return term.lower() in text.lower()

    def _proximity_match(text):
        """Return True if all search terms appear within 'proximity' words of each other."""
        words = re.findall(r'\S+', text.lower())
        term_positions = {}
        for term in search_terms:
            term_lower = term.lower()
            positions = []
            if use_regex:
                for i, word in enumerate(words):
                    if re.search(term, word, re.IGNORECASE):
                        positions.append(i)
            elif use_fuzzy:
                for i, word in enumerate(words):
                    clean = re.sub(r'^[^\w]+|[^\w]+$', '', word)
                    if clean and fuzz.ratio(term_lower, clean) >= FUZZY_THRESHOLD:
                        positions.append(i)
            else:
                for i, word in enumerate(words):
                    if term_lower in word:
                        positions.append(i)
            if not positions:
                return False
            term_positions[term] = positions
        for combo in product(*term_positions.values()):
            if max(combo) - min(combo) <= proximity:
                return True
        return False

    def _excludes_match(text):
        """Return True if any exclude term is found in text."""
        if use_regex:
            return any(re.search(t, text, re.IGNORECASE) for t in exclude_terms)
        if use_fuzzy:
            return any(_fuzzy_word_match(text, t) is not None for t in exclude_terms)
        if use_whole_word:
            return any(re.search(_whole_word_pattern(t), text, re.IGNORECASE) for t in exclude_terms)
        text_lower = text.lower()
        return any(t.lower() in text_lower for t in exclude_terms)

    def text_matches(text):
        """Return True if search terms are found in text (ANY or ALL based on mode)."""
        if expression_ast is not None:
            from peekdocs.expr_parser import evaluate_expression
            return evaluate_expression(expression_ast, text, _single_term_matches, filename=filename)
        check = all if match_all else any
        if use_regex:
            if use_proximity:
                matched = _proximity_match(text)
            else:
                matched = check(re.search(term, text, re.IGNORECASE) for term in search_terms)
        elif use_fuzzy:
            if use_proximity:
                matched = _proximity_match(text)
            else:
                matched = check(_fuzzy_word_match(text, term) is not None for term in search_terms)
        elif use_whole_word:
            matched = check(
                bool(re.search(_whole_word_pattern(term), text, re.IGNORECASE))
                for term in search_terms
            )
            if matched and use_proximity:
                matched = _proximity_match(text)
        else:
            text_lower = text.lower()
            matched = check(term.lower() in text_lower for term in search_terms)
            if matched and use_proximity:
                matched = _proximity_match(text)
        if not matched:
            return False
        if exclude_terms:
            return not _excludes_match(text)
        return True

    def highlight_text(text):
        """Apply ** highlighting around matched search terms."""
        if expression_ast is not None:
            from peekdocs.expr_parser import extract_positive_terms
            terms_to_highlight = extract_positive_terms(expression_ast)
        else:
            terms_to_highlight = search_terms
        if use_fuzzy:
            highlighted = text
            for term in terms_to_highlight:
                result_words = []
                for word in highlighted.split():
                    clean = re.sub(r'^[^\w]+|[^\w]+$', '', word)
                    if clean and fuzz.ratio(term.lower(), clean.lower()) >= FUZZY_THRESHOLD:
                        result_words.append(f"**{word}**")
                    else:
                        result_words.append(word)
                highlighted = " ".join(result_words)
            return highlighted
        highlighted = text
        for term in terms_to_highlight:
            if use_wildcard:
                pattern = _wildcard_to_regex(term)
            elif use_regex:
                pattern = term
            elif use_whole_word:
                pattern = _whole_word_pattern(term)
            else:
                pattern = re.escape(term)
            highlighted = re.sub(pattern, lambda m: f"**{m.group()}**", highlighted, flags=re.IGNORECASE)
        return highlighted

    def context_group_to_match(group, file_dir, filename):
        """Convert a context group to a match tuple with pre-highlighted text."""
        first_match_num = next(ln for ln, _, is_match in group if is_match)
        parts = []
        for ln, text, is_match in group:
            parts.append(text)
        return (file_dir, filename, first_match_num, "\n".join(parts))

    matches = []

    def _line_passes(text):
        """Return True if text passes both text matching and content range filters."""
        if range_only:
            text_ok = True
        else:
            text_ok = text_matches(text)
        if not text_ok:
            return False
        if content_ranges and not line_matches_content_ranges(text, content_ranges):
            return False
        return True

    line_proximity = config.get("line_proximity", 0)

    if use_context:
        match_indices = {i for i, (_, text) in enumerate(all_lines) if _line_passes(text)}
        if match_indices:
            groups = apply_context(all_lines, match_indices, context_before, context_after)
            for group in groups:
                matches.append(context_group_to_match(group, file_dir, filename))
    else:
        for line_num, text in all_lines:
            if _line_passes(text):
                matches.append((file_dir, filename, line_num, text))

    # Line proximity filter: keep only matches where all search terms
    # appear within N lines of each other in the same file
    if line_proximity > 0 and len(search_terms) > 1 and matches:
        # Build a map: for each term, which line numbers matched it?
        term_lines = {}
        for term in search_terms:
            term_lower = term.lower()
            term_lines[term] = set()
            for fd, fn, ln, text in matches:
                text_lower = text.lower()
                if use_regex:
                    if re.search(term, text, re.IGNORECASE):
                        term_lines[term].add(ln)
                elif use_whole_word:
                    if re.search(_whole_word_pattern(term), text, re.IGNORECASE):
                        term_lines[term].add(ln)
                elif use_fuzzy:
                    if _fuzzy_word_match(text, term) is not None:
                        term_lines[term].add(ln)
                else:
                    if term_lower in text_lower:
                        term_lines[term].add(ln)

        # Find line numbers where all terms are within line_proximity lines,
        # then group them into a single combined match with all lines together
        if all(term_lines.values()):
            from itertools import product as _lp_product
            # Collect all valid proximity windows
            windows = []
            for combo in _lp_product(*[sorted(tl) for tl in term_lines.values()]):
                if max(combo) - min(combo) <= line_proximity:
                    windows.append((min(combo), max(combo)))
            if windows:
                # Merge overlapping windows
                windows.sort()
                merged = [windows[0]]
                for start, end in windows[1:]:
                    if start <= merged[-1][1] + 1:
                        merged[-1] = (merged[-1][0], max(merged[-1][1], end))
                    else:
                        merged.append((start, end))
                # Build a lookup of all lines by line number
                all_lines_dict = {ln: text for ln, text in all_lines}
                # Create one grouped match per window
                grouped = []
                for win_start, win_end in merged:
                    lines_in_window = []
                    for ln in range(win_start, win_end + 1):
                        if ln in all_lines_dict:
                            lines_in_window.append(all_lines_dict[ln])
                    if lines_in_window:
                        combined_text = "\n".join(lines_in_window)
                        grouped.append((file_dir, filename, win_start, combined_text))
                matches = grouped
            else:
                matches = []
        else:
            matches = []

    return (matches, [])


_DEFAULT_MAX_FILE_SIZE_MB = 100


def _process_file(args_tuple):
    """Process a single file and return (matches, skipped) for that file."""
    filepath, config = args_tuple
    filename = os.path.basename(filepath)

    # File size guard — skip files over the configured limit
    max_mb = config.get("max_file_size_mb", _DEFAULT_MAX_FILE_SIZE_MB)
    if max_mb > 0:
        try:
            size_mb = os.path.getsize(filepath) / (1024 * 1024)
            if size_mb > max_mb:
                return ([], [(filename,
                    f"Skipped — file is {size_mb:.0f} MB, exceeds the {max_mb} MB limit. "
                    f"GUI: increase Max File Size in Advanced Search Options. "
                    f"CLI: use --max-file-size 0 (no limit) or --config max_file_size_mb=0 to save permanently.")])
        except OSError:
            pass

    # Metadata range filtering — skip entire file if it doesn't match
    metadata_ranges = config.get("metadata_ranges", [])
    if metadata_ranges:
        if not file_matches_metadata_ranges(filepath, metadata_ranges):
            return ([], [])

    # Filename range filtering — skip entire file if filename doesn't match
    filename_ranges = config.get("filename_ranges", [])
    if filename_ranges:
        if not file_matches_filename_ranges(filename, filename_ranges):
            return ([], [])

    try:
        all_lines = _extract_lines(
            filepath,
            config.get("use_ocr", False),
            config.get("_ocr_image_func"),
        )
    except MemoryError:
        return ([], [(filename,
            f"Skipped — '{filename}' caused an out-of-memory error. "
            f"The file may be too large to process. Try reducing CPU cores (-c 1) "
            f"or excluding this file type.")])
    except Exception as e:
        return ([], [(filename, _friendly_file_error(e, filename))])

    return _search_file_lines(
        all_lines,
        os.path.dirname(filepath),
        filename,
        config,
    )


def _friendly_file_error(exc, filename):
    """Return a user-friendly error message for file processing failures."""
    if isinstance(exc, PermissionError):
        return (f"Permission denied — '{filename}' may be open in another program. "
                "Close it and try again.")
    if isinstance(exc, OSError):
        errno = getattr(exc, 'errno', None)
        # Windows: file in use (errno 13, 32, 33)
        if errno in (13, 32, 33):
            return (f"'{filename}' is locked or in use by another program. "
                    "Close it and try again.")
    exc_str = str(exc).lower()
    # Password-protected archives
    if "password" in exc_str or "encrypted" in exc_str or "RuntimeError" in type(exc).__name__:
        ext = os.path.splitext(filename)[1].lower()
        if ext in (".zip", ".7z", ".rar", ".tar", ".gz", ".bz2", ".tgz"):
            return f"'{filename}' appears to be password-protected. peekdocs cannot read encrypted archives."
    # OneDrive placeholder files (cloud-only, not downloaded)
    if "cloud" in exc_str or "not available" in exc_str or (isinstance(exc, FileNotFoundError) and os.path.exists(filename)):
        return f"'{filename}' may be a cloud-only placeholder (e.g., OneDrive). Download the file first."
    return str(exc)


def discover_files(cwd, recursive, use_ocr, file_types=None, file_names=None):
    """Find all searchable files in cwd, applying type/name filters.

    Returns a sorted list of file paths on success, or (exit_code, message)
    on error (e.g., specified files not found).
    """
    if recursive:
        glob_prefix = os.path.join(cwd, "**", "*")
    else:
        glob_prefix = os.path.join(cwd, "*")

    # System/OS files to exclude from search results (peekdocs's own
    # files are caught by the is_peekdocs_internal_file() helper inside
    # the discovery loop below — no enumeration needed because the
    # naming convention reserves the peekdocs_ / .peekdocs prefixes).
    _EXCLUDE_NAMES = {
        "thumbs.db", "desktop.ini", ".ds_store", ".ds_store?",
        ".spotlight-v100", ".trashes", ".fseventsd",
    }

    # Filenames that are searchable despite having no standard extension
    _SPECIAL_FILENAMES = {".env", "dockerfile", ".dockerfile"}

    # Discover all supported file types dynamically
    discovered = []
    search_types = SUPPORTED_TYPES | (OCR_IMAGE_TYPES if use_ocr else set())
    for ext in sorted(search_types):
        matches = glob.glob(glob_prefix + ext, recursive=recursive)
        for f in matches:
            basename = os.path.basename(f)
            if basename.lower() in _EXCLUDE_NAMES:
                continue
            if is_peekdocs_internal_file(basename):
                continue
            if basename.startswith("~$") or basename.startswith("~"):
                continue  # Word/Excel lock files and temp files
            if basename.startswith("._"):
                continue  # macOS resource fork shadow files
            if os.path.islink(f):
                continue  # Skip symlinks to prevent infinite loops
            try:
                mode = os.stat(f).st_mode
                if stat.S_ISFIFO(mode) or stat.S_ISSOCK(mode):
                    continue  # Named pipes and sockets would hang
            except OSError:
                continue
            discovered.append(f)
    # Also discover files with special names (no standard extension).
    # Glob("*") skips dotfiles, so also glob(".*") to catch .env etc.
    _special_seen = set(discovered)
    for pattern_suffix in ("*", ".*"):
        _sp_glob = os.path.join(cwd, "**", pattern_suffix) if recursive else os.path.join(cwd, pattern_suffix)
        for f in glob.glob(_sp_glob, recursive=recursive):
            if f in _special_seen:
                continue
            basename = os.path.basename(f).lower()
            if basename in _SPECIAL_FILENAMES:
                discovered.append(f)
                _special_seen.add(f)
    # Skip virtual/system directories that could hang or return infinite data
    _SKIP_DIRS = ("/proc/", "/sys/", "/dev/", "/.gvfs/")
    discovered = [f for f in discovered if not any(skip in f for skip in _SKIP_DIRS)]

    all_files = sorted(discovered)

    if file_types is not None:
        all_files = [f for f in all_files if os.path.splitext(f)[1].lower() in file_types]

    if file_names is not None:
        name_set = {n.lower() for n in file_names}
        all_files = [f for f in all_files if os.path.basename(f).lower() in name_set]
        missing = name_set - {os.path.basename(f).lower() for f in all_files}
        if missing:
            msg = "\n".join(f"File not found: {m}" for m in sorted(missing)) + "\n"
            return (2, msg)

    return all_files
