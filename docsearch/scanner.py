"""File processing and discovery for docsearch."""

import csv
import glob
import os
import re
from copy import deepcopy
from html.parser import HTMLParser
from itertools import product

import fitz
from docx import Document
from odf.opendocument import load as load_odt
from odf.text import P as OdtParagraph
from odf.table import Table as OdfTable, TableRow as OdfTableRow, TableCell as OdfTableCell
from odf import teletype
from openpyxl import load_workbook
from striprtf.striprtf import rtf_to_text
from pptx import Presentation as PptxPresentation
import ebooklib
from ebooklib import epub

from docsearch.constants import SUPPORTED_TYPES, OCR_IMAGE_TYPES, FUZZY_THRESHOLD


def apply_context(all_lines, match_indices, before, after):
    """Expand match indices with before/after context and return merged groups."""
    if not match_indices:
        return []

    total = len(all_lines)
    ranges = []
    for idx in sorted(match_indices):
        start = max(0, idx - before)
        end = min(total - 1, idx + after)
        ranges.append((start, end))

    merged = []
    for start, end in ranges:
        if merged and start <= merged[-1][1] + 1:
            merged[-1] = (merged[-1][0], max(merged[-1][1], end))
        else:
            merged.append([start, end])

    groups = []
    for start, end in merged:
        group = []
        for i in range(start, end + 1):
            line_num, text = all_lines[i]
            is_match = i in match_indices
            group.append((line_num, text, is_match))
        groups.append(group)

    return groups


def _wildcard_to_regex(term):
    """Convert a wildcard pattern (* and ?) to a regex pattern."""
    escaped = re.escape(term)
    escaped = escaped.replace(r'\*', r'\w*')
    escaped = escaped.replace(r'\?', r'\w')
    return r'\b' + escaped + r'\b'


def _ocr_image(image):
    """Run OCR on a PIL Image and return extracted text."""
    import pytesseract
    return pytesseract.image_to_string(image)


def _extract_lines(filepath, use_ocr=False, ocr_func=None):
    """Extract text lines from a file.

    Returns a list of (line_num, text) tuples. The meaning of line_num
    varies by file type (line number, page number, row number, etc.).
    Raises an exception on error — the caller decides how to handle it.
    """
    if ocr_func is None:
        ocr_func = _ocr_image

    ext = os.path.splitext(filepath)[1].lower()
    all_lines = []

    if ext == ".docx":
        doc = Document(filepath)
        all_lines = [(i, para.text) for i, para in enumerate(doc.paragraphs, start=1)]

    elif ext == ".pdf":
        doc = fitz.open(filepath)
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            if use_ocr and len((text or "").strip()) < 10:
                import io
                from PIL import Image
                pix = page.get_pixmap(dpi=300)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = ocr_func(img)
            if not text or not text.strip():
                continue
            for line in text.split("\n"):
                all_lines.append((page_num, line))
        doc.close()

    elif ext == ".csv":
        with open(filepath, newline="", encoding="utf-8", errors="replace") as csvfile:
            reader = csv.reader(csvfile)
            all_lines = [(row_num, ", ".join(row)) for row_num, row in enumerate(reader, start=1)]

    elif ext == ".tsv":
        with open(filepath, newline="", encoding="utf-8", errors="replace") as tsvfile:
            reader = csv.reader(tsvfile, delimiter="\t")
            all_lines = [(row_num, "\t".join(row)) for row_num, row in enumerate(reader, start=1)]

    elif ext == ".epub":
        book = epub.read_epub(filepath)
        line_num = 0
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            html_content = item.get_content().decode("utf-8", errors="replace")
            class _EpubTextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text_parts = []
                def handle_data(self, data):
                    self.text_parts.append(data)
            parser = _EpubTextExtractor()
            parser.feed(html_content)
            for line in "".join(parser.text_parts).split("\n"):
                stripped = line.strip()
                if stripped:
                    line_num += 1
                    all_lines.append((line_num, stripped))

    elif ext == ".odt":
        odt_doc = load_odt(filepath)
        all_lines = [(i, teletype.extractText(para)) for i, para in enumerate(odt_doc.getElementsByType(OdtParagraph), start=1)]

    elif ext == ".odp":
        odp_doc = load_odt(filepath)
        all_lines = [(i, teletype.extractText(para)) for i, para in enumerate(odp_doc.getElementsByType(OdtParagraph), start=1)]

    elif ext == ".ods":
        ods_doc = load_odt(filepath)
        row_num = 0
        for table in ods_doc.getElementsByType(OdfTable):
            for row in table.getElementsByType(OdfTableRow):
                row_num += 1
                cells = []
                for cell in row.getElementsByType(OdfTableCell):
                    cell_text = teletype.extractText(cell)
                    if cell_text:
                        cells.append(cell_text)
                row_text = ", ".join(cells)
                all_lines.append((row_num, row_text))

    elif ext in (".txt", ".md", ".json", ".xml", ".log", ".yaml", ".yml", ".toml", ".rst", ".tex", ".ini", ".cfg", ".sql"):
        with open(filepath, encoding="utf-8", errors="replace") as txtfile:
            all_lines = [(line_num, line.rstrip("\n")) for line_num, line in enumerate(txtfile, start=1)]

    elif ext == ".html":
        class _HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text_parts = []
            def handle_data(self, data):
                self.text_parts.append(data)
        with open(filepath, encoding="utf-8", errors="replace") as htmlfile:
            parser = _HTMLTextExtractor()
            parser.feed(htmlfile.read())
        lines = "".join(parser.text_parts).split("\n")
        all_lines = [(line_num, line.strip()) for line_num, line in enumerate(lines, start=1)]

    elif ext == ".pptx":
        prs = PptxPresentation(filepath)
        para_num = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_num += 1
                        all_lines.append((para_num, para.text))

    elif ext == ".rtf":
        with open(filepath, encoding="utf-8", errors="replace") as rtffile:
            raw = rtffile.read()
        plain = rtf_to_text(raw)
        all_lines = [(line_num, line) for line_num, line in enumerate(plain.split("\n"), start=1)]

    elif ext == ".xlsx":
        wb = load_workbook(filepath, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            for row_num, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                row_text = ", ".join(str(cell) for cell in row if cell is not None)
                all_lines.append((row_num, row_text))
        wb.close()

    elif ext in OCR_IMAGE_TYPES:
        from PIL import Image
        img = Image.open(filepath)
        text = ocr_func(img)
        if text and text.strip():
            all_lines = [(line_num, line) for line_num, line in enumerate(text.split("\n"), start=1) if line.strip()]

    return all_lines


def _search_file_lines(all_lines, file_dir, filename, config):
    """Search extracted lines for matches.

    Args:
        all_lines: list of (line_num, text) tuples from _extract_lines
        file_dir: directory containing the file
        filename: base name of the file
        config: search configuration dict

    Returns:
        (matches, skipped) where:
        - matches = list of (file_dir, filename, line_num, text) tuples
        - skipped = list of (filename, error_msg) tuples (usually empty)
    """
    search_terms = config["search_terms"]
    use_regex = config["use_regex"]
    match_all = config["match_all"]
    use_proximity = config["use_proximity"]
    proximity = config["proximity"]
    use_context = config["use_context"]
    context_before = config["context_before"]
    context_after = config["context_after"]
    use_fuzzy = config.get("use_fuzzy", False)
    exclude_terms = config.get("exclude_terms", [])
    use_wildcard = config.get("use_wildcard", False)
    if use_wildcard:
        search_terms = [_wildcard_to_regex(t) for t in search_terms]
        if exclude_terms:
            exclude_terms = [_wildcard_to_regex(t) for t in exclude_terms]
        use_regex = True
    if use_fuzzy:
        from rapidfuzz import fuzz

    def _fuzzy_word_match(text, term):
        """Return the first word in text that fuzzy-matches term, or None."""
        for word in re.findall(r'\S+', text):
            clean = re.sub(r'^[^\w]+|[^\w]+$', '', word)
            if not clean:
                continue
            if fuzz.ratio(term.lower(), clean.lower()) >= FUZZY_THRESHOLD:
                return word
        return None

    def _proximity_match(text):
        """Return True if all search terms appear within 'proximity' words of each other."""
        words = re.findall(r'\S+', text.lower())
        term_positions = {}
        for term in search_terms:
            term_lower = term.lower()
            positions = []
            if use_regex:
                for i, word in enumerate(words):
                    if re.search(term, word, re.IGNORECASE):
                        positions.append(i)
            elif use_fuzzy:
                for i, word in enumerate(words):
                    clean = re.sub(r'^[^\w]+|[^\w]+$', '', word)
                    if clean and fuzz.ratio(term_lower, clean) >= FUZZY_THRESHOLD:
                        positions.append(i)
            else:
                for i, word in enumerate(words):
                    if term_lower in word:
                        positions.append(i)
            if not positions:
                return False
            term_positions[term] = positions
        for combo in product(*term_positions.values()):
            if max(combo) - min(combo) <= proximity:
                return True
        return False

    def _excludes_match(text):
        """Return True if any exclude term is found in text."""
        if use_regex:
            return any(re.search(t, text, re.IGNORECASE) for t in exclude_terms)
        if use_fuzzy:
            return any(_fuzzy_word_match(text, t) is not None for t in exclude_terms)
        text_lower = text.lower()
        return any(t.lower() in text_lower for t in exclude_terms)

    def text_matches(text):
        """Return True if search terms are found in text (ANY or ALL based on mode)."""
        check = all if match_all else any
        if use_regex:
            if use_proximity:
                matched = _proximity_match(text)
            else:
                matched = check(re.search(term, text, re.IGNORECASE) for term in search_terms)
        elif use_fuzzy:
            if use_proximity:
                matched = _proximity_match(text)
            else:
                matched = check(_fuzzy_word_match(text, term) is not None for term in search_terms)
        else:
            text_lower = text.lower()
            matched = check(term.lower() in text_lower for term in search_terms)
            if matched and use_proximity:
                matched = _proximity_match(text)
        if not matched:
            return False
        if exclude_terms:
            return not _excludes_match(text)
        return True

    def highlight_text(text):
        """Apply ** highlighting around matched search terms."""
        if use_fuzzy:
            highlighted = text
            for term in search_terms:
                result_words = []
                for word in highlighted.split():
                    clean = re.sub(r'^[^\w]+|[^\w]+$', '', word)
                    if clean and fuzz.ratio(term.lower(), clean.lower()) >= FUZZY_THRESHOLD:
                        result_words.append(f"**{word}**")
                    else:
                        result_words.append(word)
                highlighted = " ".join(result_words)
            return highlighted
        highlighted = text
        for term in search_terms:
            pattern = term if use_regex else re.escape(term)
            highlighted = re.sub(pattern, lambda m: f"**{m.group()}**", highlighted, flags=re.IGNORECASE)
        return highlighted

    def context_group_to_match(group, file_dir, filename):
        """Convert a context group to a match tuple with pre-highlighted text."""
        first_match_num = next(ln for ln, _, is_match in group if is_match)
        parts = []
        for ln, text, is_match in group:
            if is_match:
                parts.append(highlight_text(text))
            else:
                parts.append(text)
        return (file_dir, filename, first_match_num, "\n".join(parts))

    matches = []

    if use_context:
        match_indices = {i for i, (_, text) in enumerate(all_lines) if text_matches(text)}
        if match_indices:
            groups = apply_context(all_lines, match_indices, context_before, context_after)
            for group in groups:
                matches.append(context_group_to_match(group, file_dir, filename))
    else:
        for line_num, text in all_lines:
            if text_matches(text):
                if use_fuzzy:
                    matches.append((file_dir, filename, line_num, highlight_text(text)))
                else:
                    matches.append((file_dir, filename, line_num, text))

    return (matches, [])


def _process_file(args_tuple):
    """Process a single file and return (matches, skipped) for that file."""
    filepath, config = args_tuple
    filename = os.path.basename(filepath)

    try:
        all_lines = _extract_lines(
            filepath,
            config.get("use_ocr", False),
            config.get("_ocr_image_func"),
        )
    except Exception as e:
        return ([], [(filename, _friendly_file_error(e, filename))])

    return _search_file_lines(
        all_lines,
        os.path.dirname(filepath),
        filename,
        config,
    )


def _friendly_file_error(exc, filename):
    """Return a user-friendly error message for file processing failures."""
    if isinstance(exc, PermissionError):
        return (f"Permission denied — '{filename}' may be open in another program. "
                "Close it and try again.")
    if isinstance(exc, OSError):
        errno = getattr(exc, 'errno', None)
        # Windows: file in use (errno 13, 32, 33)
        if errno in (13, 32, 33):
            return (f"'{filename}' is locked or in use by another program. "
                    "Close it and try again.")
    return str(exc)


def discover_files(cwd, recursive, use_ocr, file_types=None, file_names=None):
    """Find all searchable files in cwd, applying type/name filters.

    Returns a sorted list of file paths on success, or (exit_code, message)
    on error (e.g., specified files not found).
    """
    if recursive:
        glob_prefix = os.path.join(cwd, "**", "*")
    else:
        glob_prefix = os.path.join(cwd, "*")

    docx_files = sorted(
        f for f in glob.glob(glob_prefix + ".docx", recursive=recursive)
        if os.path.basename(f) != "docsearch_results.docx"
        and not os.path.basename(f).startswith("DO_NOT_SEARCH_")
    )
    pdf_files = sorted(glob.glob(glob_prefix + ".pdf", recursive=recursive))
    csv_files = sorted(
        f for f in glob.glob(glob_prefix + ".csv", recursive=recursive)
        if os.path.basename(f) != "docsearch_results.csv"
        and not os.path.basename(f).startswith("DO_NOT_SEARCH_")
    )
    odt_files = sorted(glob.glob(glob_prefix + ".odt", recursive=recursive))
    txt_files = sorted(
        f for f in glob.glob(glob_prefix + ".txt", recursive=recursive)
        if os.path.basename(f) != "docsearch_results.txt"
        and not os.path.basename(f).startswith("DO_NOT_SEARCH_")
    )
    html_files = sorted(glob.glob(glob_prefix + ".html", recursive=recursive))
    xlsx_files = sorted(glob.glob(glob_prefix + ".xlsx", recursive=recursive))
    md_files = sorted(glob.glob(glob_prefix + ".md", recursive=recursive))
    json_files = sorted(
        f for f in glob.glob(glob_prefix + ".json", recursive=recursive)
        if os.path.basename(f) != "docsearch_results.json"
        and not os.path.basename(f).startswith("DO_NOT_SEARCH_")
    )
    rtf_files = sorted(glob.glob(glob_prefix + ".rtf", recursive=recursive))
    pptx_files = sorted(glob.glob(glob_prefix + ".pptx", recursive=recursive))
    xml_files = sorted(glob.glob(glob_prefix + ".xml", recursive=recursive))
    log_files = sorted(
        f for f in glob.glob(glob_prefix + ".log", recursive=recursive)
        if os.path.basename(f) != "docsearch_errors.log"
    )
    yaml_files = sorted(glob.glob(glob_prefix + ".yaml", recursive=recursive))
    yml_files = sorted(glob.glob(glob_prefix + ".yml", recursive=recursive))
    tsv_files = sorted(glob.glob(glob_prefix + ".tsv", recursive=recursive))
    epub_files = sorted(glob.glob(glob_prefix + ".epub", recursive=recursive))
    ods_files = sorted(glob.glob(glob_prefix + ".ods", recursive=recursive))
    odp_files = sorted(glob.glob(glob_prefix + ".odp", recursive=recursive))
    toml_files = sorted(glob.glob(glob_prefix + ".toml", recursive=recursive))
    rst_files = sorted(glob.glob(glob_prefix + ".rst", recursive=recursive))
    tex_files = sorted(glob.glob(glob_prefix + ".tex", recursive=recursive))
    ini_files = sorted(glob.glob(glob_prefix + ".ini", recursive=recursive))
    cfg_files = sorted(glob.glob(glob_prefix + ".cfg", recursive=recursive))
    sql_files = sorted(glob.glob(glob_prefix + ".sql", recursive=recursive))
    if use_ocr:
        jpg_files = sorted(glob.glob(glob_prefix + ".jpg", recursive=recursive))
        jpeg_files = sorted(glob.glob(glob_prefix + ".jpeg", recursive=recursive))
        png_files = sorted(glob.glob(glob_prefix + ".png", recursive=recursive))
        tiff_files = sorted(glob.glob(glob_prefix + ".tiff", recursive=recursive))
        tif_files = sorted(glob.glob(glob_prefix + ".tif", recursive=recursive))
        bmp_files = sorted(glob.glob(glob_prefix + ".bmp", recursive=recursive))
        image_files = jpg_files + jpeg_files + png_files + tiff_files + tif_files + bmp_files
    else:
        image_files = []
    all_files = sorted(
        f for f in docx_files + pdf_files + csv_files + odt_files + txt_files + html_files + xlsx_files + md_files + json_files + rtf_files + pptx_files + xml_files + log_files + yaml_files + yml_files + tsv_files + epub_files + ods_files + odp_files + toml_files + rst_files + tex_files + ini_files + cfg_files + sql_files + image_files
        if not os.path.basename(f).startswith("DO_NOT_SEARCH")
    )

    if file_types is not None:
        all_files = [f for f in all_files if os.path.splitext(f)[1].lower() in file_types]

    if file_names is not None:
        name_set = {n.lower() for n in file_names}
        all_files = [f for f in all_files if os.path.basename(f).lower() in name_set]
        missing = name_set - {os.path.basename(f).lower() for f in all_files}
        if missing:
            msg = "\n".join(f"File not found: {m}" for m in sorted(missing)) + "\n"
            return (2, msg)

    return all_files
