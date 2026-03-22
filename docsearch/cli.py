"""Command-line interface for DocSearch."""

import csv
from copy import deepcopy
import glob
from html.parser import HTMLParser
from itertools import product
import logging
import multiprocessing
import threading
import os
import re
import shutil
import sys
import textwrap
import time
from datetime import datetime

logging.getLogger("pymupdf").setLevel(logging.ERROR)

import warnings
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")

import fitz
from docx import Document
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from odf.opendocument import load as load_odt
from odf.text import P as OdtParagraph
from odf.table import Table as OdfTable, TableRow as OdfTableRow, TableCell as OdfTableCell
from odf import teletype
from docx.enum.text import WD_COLOR_INDEX
from openpyxl import load_workbook
from striprtf.striprtf import rtf_to_text
from pptx import Presentation as PptxPresentation
import ebooklib
from ebooklib import epub
from importlib.metadata import version as pkg_version


VERSION = pkg_version("claude-docsearch")

SUPPORTED_TYPES = {".docx", ".pdf", ".csv", ".odt", ".txt", ".html", ".xlsx", ".md", ".json", ".rtf", ".pptx", ".xml", ".log", ".yaml", ".yml", ".tsv", ".epub", ".ods", ".odp", ".toml", ".rst", ".tex", ".ini", ".cfg", ".sql"}

BANNER_TOP = (
    '\n OR search — finds paragraphs containing ANY of the search terms. Example: docsearch term1 term2 term3\n'
    'AND search — finds paragraphs containing ALL of the search terms. Example: docsearch -a term1 term2 term3\n'
    'Use option flag -a for AND searches. Example: docsearch -a term1 term2 term3\n'
    'Use option flag -A to show lines after each match. Example: docsearch -A 5 term1\n'
    'Use option flag -B to show lines before each match. Example: docsearch -B 5 term1\n'
    'Use option flag -c to set number of CPU cores. Example: docsearch -c 4 budget revenue\n'
    'Use option flag -f to search specific files. Example: docsearch -f report.pdf,notes.txt term1\n'
    'Use option flag -h for help. Example: docsearch -h     (Also displays common Regex patterns)\n'
    'Use option flag -p to find terms within N words of each other. Example: docsearch -p 5 budget revenue'
)

BANNER_BOTTOM = (
    'Use option flag -r to search subdirectories. Example: docsearch -r term1 term2 term3\n'
    'Use option flag -s to save the last search report. Example: docsearch -s name_of_my_file\n'
    'Use option flag -sa to search and auto-append results to a named file. Example: docsearch -sa my_report budget revenue\n'
    'Use option flag -t to filter by file type. Example: docsearch -t pdf,docx term1 term2\n'
    'Use option flag -v for version. Example: docsearch -v\n'
    'Use option flag -x for regex searches. Example: docsearch -x "\\d{3}-\\d{3}-\\d{4}"\n'
    'Special characters (<, >, [, ], *, ?, $, |, etc.) must be enclosed in quotes\n'
    'More details here: https://github.com/exbuf/Claude-DocSearch/blob/main/README.md'
)

REGEX_PATTERNS = (
    '\nCommon Regex Search Patterns (enclose in quotes):\n'
    '  \\d{3}-\\d{3}-\\d{4}                              US phone numbers (555-123-4567)\n'
    '  [A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\\.[A-Z]{2,}    Email addresses (jane@example.com)\n'
    '  \\d{4}-\\d{2}-\\d{2}                              Dates, YYYY-MM-DD (2026-03-17)\n'
    '  \\$\\d+(\\.\\d{2})?                                 Dollar amounts ($45.99)\n'
    '  \\d{3}-\\d{2}-\\d{4}                              SSN format (123-45-6789)\n'
    '  \\d{1,3}\\.\\d{1,3}\\.\\d{1,3}\\.\\d{1,3}              IP addresses (192.168.1.1)\n'
    '  https?://\\S+                                    URLs (https://example.com)\n'
    '  \\b[A-Z]{2,}\\b                                   Acronyms, all caps (NASA, FBI)\n'
    '  \\b\\d{5}(-\\d{4})?\\b                              US ZIP codes (12345 or 12345-6789)\n'
    '  \\(\\d{3}\\)\\s?\\d{3}-\\d{4}                         Phone with area code parens ((555) 123-4567)\n'
    '  \\b[A-Z][a-z]+\\s[A-Z][a-z]+\\b                    Proper names (John Smith)\n'
    '  \\b\\d+%                                          Percentages (92%)\n'
    '  Q[1-4]\\s?\\d{4}                                  Fiscal quarters (Q1 2026)\n'
)


def apply_context(all_lines, match_indices, before, after):
    """Expand match indices with before/after context and return merged groups."""
    if not match_indices:
        return []

    total = len(all_lines)
    ranges = []
    for idx in sorted(match_indices):
        start = max(0, idx - before)
        end = min(total - 1, idx + after)
        ranges.append((start, end))

    merged = []
    for start, end in ranges:
        if merged and start <= merged[-1][1] + 1:
            merged[-1] = (merged[-1][0], max(merged[-1][1], end))
        else:
            merged.append([start, end])

    groups = []
    for start, end in merged:
        group = []
        for i in range(start, end + 1):
            line_num, text = all_lines[i]
            is_match = i in match_indices
            group.append((line_num, text, is_match))
        groups.append(group)

    return groups


def _default_cores():
    """Return the default number of worker processes: half of available cores, minimum 1."""
    cpu = os.cpu_count()
    if cpu is None:
        return 1
    return max(1, cpu // 2)


def _load_config():
    """Load defaults from ~/.docsearchrc if it exists."""
    config_path = os.path.join(os.path.expanduser("~"), ".docsearchrc")
    if not os.path.exists(config_path):
        return {}
    config = {}
    bool_keys = {"recursive", "quiet", "match_all", "regex"}
    int_keys = {"cores", "context_before", "context_after"}
    str_keys = {"file_types"}
    with open(config_path) as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if "=" not in line:
                continue
            key, _, value = line.partition("=")
            key = key.strip()
            value = value.strip()
            if key in bool_keys:
                config[key] = value.lower() in ("true", "yes", "1")
            elif key in int_keys:
                try:
                    config[key] = int(value)
                except ValueError:
                    pass
            elif key in str_keys:
                if value:
                    config[key] = value
    return config


def _process_file(args_tuple):
    """Process a single file and return (matches, skipped) for that file."""
    filepath, config = args_tuple

    search_terms = config["search_terms"]
    use_regex = config["use_regex"]
    match_all = config["match_all"]
    use_proximity = config["use_proximity"]
    proximity = config["proximity"]
    use_context = config["use_context"]
    context_before = config["context_before"]
    context_after = config["context_after"]

    def _proximity_match(text):
        """Return True if all search terms appear within 'proximity' words of each other."""
        words = re.findall(r'\S+', text.lower())
        term_positions = {}
        for term in search_terms:
            term_lower = term.lower()
            positions = []
            if use_regex:
                for i, word in enumerate(words):
                    if re.search(term, word, re.IGNORECASE):
                        positions.append(i)
            else:
                for i, word in enumerate(words):
                    if term_lower in word:
                        positions.append(i)
            if not positions:
                return False
            term_positions[term] = positions
        for combo in product(*term_positions.values()):
            if max(combo) - min(combo) <= proximity:
                return True
        return False

    def text_matches(text):
        """Return True if search terms are found in text (ANY or ALL based on mode)."""
        check = all if match_all else any
        if use_regex:
            if use_proximity:
                return _proximity_match(text)
            return check(re.search(term, text, re.IGNORECASE) for term in search_terms)
        text_lower = text.lower()
        if not check(term.lower() in text_lower for term in search_terms):
            return False
        if use_proximity:
            return _proximity_match(text)
        return True

    def highlight_text(text):
        """Apply ** highlighting around matched search terms."""
        highlighted = text
        for term in search_terms:
            pattern = term if use_regex else re.escape(term)
            highlighted = re.sub(pattern, lambda m: f"**{m.group()}**", highlighted, flags=re.IGNORECASE)
        return highlighted

    def context_group_to_match(group, file_dir, filename, line_num_override=None):
        """Convert a context group to a match tuple with pre-highlighted text."""
        first_match_num = line_num_override or next(ln for ln, _, is_match in group if is_match)
        parts = []
        for ln, text, is_match in group:
            if is_match:
                parts.append(highlight_text(text))
            else:
                parts.append(text)
        return (file_dir, filename, first_match_num, "\n".join(parts))

    matches = []
    skipped = []

    def collect_matches(all_lines, file_dir, filename):
        """Find matches in all_lines and append to matches list, with context if active."""
        if use_context:
            match_indices = {i for i, (_, text) in enumerate(all_lines) if text_matches(text)}
            if not match_indices:
                return
            groups = apply_context(all_lines, match_indices, context_before, context_after)
            for group in groups:
                matches.append(context_group_to_match(group, file_dir, filename))
        else:
            for line_num, text in all_lines:
                if text_matches(text):
                    matches.append((file_dir, filename, line_num, text))

    file_dir = os.path.dirname(filepath)
    filename = os.path.basename(filepath)
    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext == ".docx":
            doc = Document(filepath)
            all_lines = [(i, para.text) for i, para in enumerate(doc.paragraphs, start=1)]
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".pdf":
            doc = fitz.open(filepath)
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text()
                if not text:
                    continue
                page_lines = [(line_num, line) for line_num, line in enumerate(text.split("\n"), start=1)]
                if use_context:
                    match_indices = {i for i, (_, lt) in enumerate(page_lines) if text_matches(lt)}
                    if match_indices:
                        groups = apply_context(page_lines, match_indices, context_before, context_after)
                        for group in groups:
                            matches.append(context_group_to_match(group, file_dir, filename, line_num_override=page_num))
                else:
                    for line_num, line in page_lines:
                        if text_matches(line):
                            matches.append((file_dir, filename, page_num, line))
            doc.close()

        elif ext == ".csv":
            with open(filepath, newline="", encoding="utf-8", errors="replace") as csvfile:
                reader = csv.reader(csvfile)
                all_lines = [(row_num, ", ".join(row)) for row_num, row in enumerate(reader, start=1)]
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".tsv":
            with open(filepath, newline="", encoding="utf-8", errors="replace") as tsvfile:
                reader = csv.reader(tsvfile, delimiter="\t")
                all_lines = [(row_num, "\t".join(row)) for row_num, row in enumerate(reader, start=1)]
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".epub":
            book = epub.read_epub(filepath)
            all_lines = []
            line_num = 0
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                html_content = item.get_content().decode("utf-8", errors="replace")
                class _EpubTextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text_parts = []
                    def handle_data(self, data):
                        self.text_parts.append(data)
                parser = _EpubTextExtractor()
                parser.feed(html_content)
                for line in "".join(parser.text_parts).split("\n"):
                    stripped = line.strip()
                    if stripped:
                        line_num += 1
                        all_lines.append((line_num, stripped))
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".odt":
            odt_doc = load_odt(filepath)
            all_lines = [(i, teletype.extractText(para)) for i, para in enumerate(odt_doc.getElementsByType(OdtParagraph), start=1)]
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".odp":
            odp_doc = load_odt(filepath)
            all_lines = [(i, teletype.extractText(para)) for i, para in enumerate(odp_doc.getElementsByType(OdtParagraph), start=1)]
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".ods":
            ods_doc = load_odt(filepath)
            all_lines = []
            row_num = 0
            for table in ods_doc.getElementsByType(OdfTable):
                for row in table.getElementsByType(OdfTableRow):
                    row_num += 1
                    cells = []
                    for cell in row.getElementsByType(OdfTableCell):
                        cell_text = teletype.extractText(cell)
                        if cell_text:
                            cells.append(cell_text)
                    row_text = ", ".join(cells)
                    all_lines.append((row_num, row_text))
            collect_matches(all_lines, file_dir, filename)

        elif ext in (".txt", ".md", ".json", ".xml", ".log", ".yaml", ".yml", ".toml", ".rst", ".tex", ".ini", ".cfg", ".sql"):
            with open(filepath, encoding="utf-8", errors="replace") as txtfile:
                all_lines = [(line_num, line.rstrip("\n")) for line_num, line in enumerate(txtfile, start=1)]
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".html":
            class _HTMLTextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text_parts = []
                def handle_data(self, data):
                    self.text_parts.append(data)
            with open(filepath, encoding="utf-8", errors="replace") as htmlfile:
                parser = _HTMLTextExtractor()
                parser.feed(htmlfile.read())
            lines = "".join(parser.text_parts).split("\n")
            all_lines = [(line_num, line.strip()) for line_num, line in enumerate(lines, start=1)]
            if use_context:
                collect_matches(all_lines, file_dir, filename)
            else:
                for line_num, text in all_lines:
                    if text and text_matches(text):
                        matches.append((file_dir, filename, line_num, text))

        elif ext == ".pptx":
            prs = PptxPresentation(filepath)
            all_lines = []
            para_num = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_num += 1
                            all_lines.append((para_num, para.text))
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".rtf":
            with open(filepath, encoding="utf-8", errors="replace") as rtffile:
                raw = rtffile.read()
            plain = rtf_to_text(raw)
            all_lines = [(line_num, line) for line_num, line in enumerate(plain.split("\n"), start=1)]
            collect_matches(all_lines, file_dir, filename)

        elif ext == ".xlsx":
            wb = load_workbook(filepath, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                sheet_lines = []
                for row_num, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    row_text = ", ".join(str(cell) for cell in row if cell is not None)
                    sheet_lines.append((row_num, row_text))
                if use_context:
                    match_indices = {i for i, (_, text) in enumerate(sheet_lines) if text and text_matches(text)}
                    if match_indices:
                        groups = apply_context(sheet_lines, match_indices, context_before, context_after)
                        for group in groups:
                            matches.append(context_group_to_match(group, file_dir, filename))
                else:
                    for row_num, row_text in sheet_lines:
                        if row_text and text_matches(row_text):
                            matches.append((file_dir, filename, row_num, row_text))
            wb.close()

    except Exception as e:
        skipped.append((filename, str(e)))

    return (matches, skipped)


def main(argv=None):
    if argv is None:
        args = sys.argv[1:]
    else:
        args = list(argv)

    original_args = list(args)

    config = _load_config()

    quiet = "-q" in args or config.get("quiet", False)
    if "-q" in args:
        args.remove("-q")

    cpu_count = os.cpu_count() or 1
    is_help = args and args[0] in ("-h", "-help", "--help")
    if not quiet:
        print(BANNER_TOP)
        if is_help:
            print('Use option flag -q to suppress the output banner. Example: docsearch -q budget revenue')
        else:
            print('Use option flag -q to suppress this output banner. Example: docsearch -q budget revenue')
        print(BANNER_BOTTOM)
        print(f'Your system has {cpu_count} CPU cores (default for -c: {max(1, cpu_count // 2)})')
        print()

    if args and args[0] in ("-v", "-version"):
        print(f"docsearch {VERSION}\n")
        return 0

    if is_help:
        if quiet:
            print(BANNER_TOP)
            print('Use option flag -q to suppress the output banner. Example: docsearch -q budget revenue')
            print(BANNER_BOTTOM)
            print(f'Your system has {cpu_count} CPU cores (default for -c: {max(1, cpu_count // 2)})')
            print()
        print(REGEX_PATTERNS)
        return 0

    if not args:
        return 0

    if args and args[0] in ("-s", "-save"):
        if len(args) < 2:
            print("No filename provided. Usage: docsearch -s name_of_your_file\n")
            return 2
        name = "_".join(args[1:]).replace(" ", "_")
        cwd = os.getcwd()
        src_docx = os.path.join(cwd, "docsearch_results.docx")
        src_txt = os.path.join(cwd, "docsearch_results.txt")
        dest_docx = os.path.join(cwd, f"DO_NOT_SEARCH_{name}.docx")
        dest_txt = os.path.join(cwd, f"DO_NOT_SEARCH_{name}.txt")
        if not os.path.exists(src_docx) or not os.path.exists(src_txt):
            print("No search results found. Run a search first.\n")
            return 2
        shutil.copy2(src_docx, dest_docx)
        shutil.copy2(src_txt, dest_txt)
        print(f"Results saved to {os.path.basename(dest_docx)} and {os.path.basename(dest_txt)}\n")
        return 0

    match_all = "-a" in args or "--all" in args or config.get("match_all", False)
    recursive = "-r" in args or config.get("recursive", False)
    use_regex = "-x" in args or config.get("regex", False)

    file_types = None
    if "-t" in args:
        idx = args.index("-t")
        if idx + 1 >= len(args):
            print("No file types provided. Usage: docsearch -t pdf,docx search_term\n")
            return 2
        raw_types = args[idx + 1].split(",")
        file_types = set()
        for t in raw_types:
            ext = "." + t.strip().lower().lstrip(".")
            if ext not in SUPPORTED_TYPES:
                print(f"Unsupported file type: {t.strip()}. Supported types: docx, pdf, csv, odt, txt, html, xlsx, md, json, rtf, pptx, xml, log, yaml, yml, tsv, epub, ods, odp, toml, rst, tex, ini, cfg, sql\n")
                return 2
            file_types.add(ext)
        args = args[:idx] + args[idx + 2:]
    elif "file_types" in config:
        raw_types = config["file_types"].split(",")
        file_types = set()
        for t in raw_types:
            ext = "." + t.strip().lower().lstrip(".")
            if ext in SUPPORTED_TYPES:
                file_types.add(ext)

    file_names = None
    if "-f" in args:
        idx = args.index("-f")
        if idx + 1 >= len(args):
            print("No file names provided. Usage: docsearch -f report.pdf,notes.txt search_term\n")
            return 2
        file_names = [n.strip() for n in args[idx + 1].split(",")]
        for n in file_names:
            ext = os.path.splitext(n)[1].lower()
            if ext not in SUPPORTED_TYPES:
                print(f"Unsupported file type in '{n}'. Supported types: docx, pdf, csv, odt, txt, html, xlsx, md, json, rtf, pptx, xml, log, yaml, yml, tsv, epub, ods, odp, toml, rst, tex, ini, cfg, sql\n")
                return 2
        args = args[:idx] + args[idx + 2:]

    if file_types is not None and file_names is not None:
        print("Cannot use -f and -t together. Use -f to search specific files or -t to filter by file type.\n")
        return 2

    context_before = config.get("context_before", 0)
    if "-B" in args:
        idx = args.index("-B")
        if idx + 1 >= len(args):
            print("No count provided. Usage: docsearch -B 5 search_term\n")
            return 2
        try:
            context_before = int(args[idx + 1])
            if context_before < 0:
                raise ValueError
        except ValueError:
            print(f"Invalid count for -B: {args[idx + 1]}. Must be a positive integer.\n")
            return 2
        args = args[:idx] + args[idx + 2:]

    context_after = config.get("context_after", 0)
    if "-A" in args:
        idx = args.index("-A")
        if idx + 1 >= len(args):
            print("No count provided. Usage: docsearch -A 5 search_term\n")
            return 2
        try:
            context_after = int(args[idx + 1])
            if context_after < 0:
                raise ValueError
        except ValueError:
            print(f"Invalid count for -A: {args[idx + 1]}. Must be a positive integer.\n")
            return 2
        args = args[:idx] + args[idx + 2:]

    proximity = 0
    if "-p" in args:
        idx = args.index("-p")
        if idx + 1 >= len(args):
            print("No count provided. Usage: docsearch -p 5 budget revenue\n")
            return 2
        try:
            proximity = int(args[idx + 1])
            if proximity < 1:
                raise ValueError
        except ValueError:
            print(f"Invalid count for -p: {args[idx + 1]}. Must be a positive integer.\n")
            return 2
        args = args[:idx] + args[idx + 2:]

    use_proximity = proximity > 0
    if use_proximity:
        match_all = True

    append_name = None
    if "-sa" in args:
        idx = args.index("-sa")
        if idx + 1 >= len(args):
            print("No filename provided. Usage: docsearch -sa my_report budget revenue\n")
            return 2
        append_name = args[idx + 1]
        args = args[:idx] + args[idx + 2:]

    cores = config.get("cores", _default_cores())
    if "-c" in args:
        idx = args.index("-c")
        if idx + 1 >= len(args):
            print("No count provided. Usage: docsearch -c 4 search_term\n")
            return 2
        try:
            cores = int(args[idx + 1])
            if cores < 1:
                raise ValueError
        except ValueError:
            print(f"Invalid count for -c: {args[idx + 1]}. Must be a positive integer.\n")
            return 2
        args = args[:idx] + args[idx + 2:]

    use_context = context_before > 0 or context_after > 0

    search_terms = [a for a in args if a not in ("-a", "--all", "-r", "-x")]

    if not search_terms:
        print("No search terms provided.\n")
        return 2

    if use_proximity and len(search_terms) < 2:
        print("Proximity search (-p) requires at least 2 search terms.\n")
        return 2

    if use_regex:
        for term in search_terms:
            try:
                re.compile(term)
            except re.error as e:
                print(f"Invalid regex pattern '{term}': {e}\n")
                return 2

    if use_regex and match_all:
        mode = "REGEX+AND"
    elif use_regex:
        mode = "REGEX"
    elif match_all:
        mode = "AND"
    else:
        mode = "OR"
    command_str = "docsearch " + " ".join(f'"{a}"' if " " in a else a for a in original_args)
    print(command_str)
    print(f"Searching ({mode}) on [{', '.join(search_terms)}] ...")
    start_time = time.time()
    cwd = os.getcwd()

    if recursive:
        glob_prefix = os.path.join(cwd, "**", "*")
    else:
        glob_prefix = os.path.join(cwd, "*")

    docx_files = sorted(
        f for f in glob.glob(glob_prefix + ".docx", recursive=recursive)
        if os.path.basename(f) != "docsearch_results.docx"
        and not os.path.basename(f).startswith("DO_NOT_SEARCH_")
    )
    pdf_files = sorted(glob.glob(glob_prefix + ".pdf", recursive=recursive))
    csv_files = sorted(glob.glob(glob_prefix + ".csv", recursive=recursive))
    odt_files = sorted(glob.glob(glob_prefix + ".odt", recursive=recursive))
    txt_files = sorted(
        f for f in glob.glob(glob_prefix + ".txt", recursive=recursive)
        if os.path.basename(f) != "docsearch_results.txt"
        and not os.path.basename(f).startswith("DO_NOT_SEARCH_")
    )
    html_files = sorted(glob.glob(glob_prefix + ".html", recursive=recursive))
    xlsx_files = sorted(glob.glob(glob_prefix + ".xlsx", recursive=recursive))
    md_files = sorted(glob.glob(glob_prefix + ".md", recursive=recursive))
    json_files = sorted(glob.glob(glob_prefix + ".json", recursive=recursive))
    rtf_files = sorted(glob.glob(glob_prefix + ".rtf", recursive=recursive))
    pptx_files = sorted(glob.glob(glob_prefix + ".pptx", recursive=recursive))
    xml_files = sorted(glob.glob(glob_prefix + ".xml", recursive=recursive))
    log_files = sorted(
        f for f in glob.glob(glob_prefix + ".log", recursive=recursive)
        if os.path.basename(f) != "docsearch_errors.log"
    )
    yaml_files = sorted(glob.glob(glob_prefix + ".yaml", recursive=recursive))
    yml_files = sorted(glob.glob(glob_prefix + ".yml", recursive=recursive))
    tsv_files = sorted(glob.glob(glob_prefix + ".tsv", recursive=recursive))
    epub_files = sorted(glob.glob(glob_prefix + ".epub", recursive=recursive))
    ods_files = sorted(glob.glob(glob_prefix + ".ods", recursive=recursive))
    odp_files = sorted(glob.glob(glob_prefix + ".odp", recursive=recursive))
    toml_files = sorted(glob.glob(glob_prefix + ".toml", recursive=recursive))
    rst_files = sorted(glob.glob(glob_prefix + ".rst", recursive=recursive))
    tex_files = sorted(glob.glob(glob_prefix + ".tex", recursive=recursive))
    ini_files = sorted(glob.glob(glob_prefix + ".ini", recursive=recursive))
    cfg_files = sorted(glob.glob(glob_prefix + ".cfg", recursive=recursive))
    sql_files = sorted(glob.glob(glob_prefix + ".sql", recursive=recursive))
    all_files = sorted(
        f for f in docx_files + pdf_files + csv_files + odt_files + txt_files + html_files + xlsx_files + md_files + json_files + rtf_files + pptx_files + xml_files + log_files + yaml_files + yml_files + tsv_files + epub_files + ods_files + odp_files + toml_files + rst_files + tex_files + ini_files + cfg_files + sql_files
        if not os.path.basename(f).startswith("DO_NOT_SEARCH")
    )

    if file_types is not None:
        all_files = [f for f in all_files if os.path.splitext(f)[1].lower() in file_types]

    if file_names is not None:
        name_set = {n.lower() for n in file_names}
        all_files = [f for f in all_files if os.path.basename(f).lower() in name_set]
        missing = name_set - {os.path.basename(f).lower() for f in all_files}
        if missing:
            for m in sorted(missing):
                print(f"File not found: {m}")
            print()
            return 2

    search_config = {
        "search_terms": search_terms,
        "use_regex": use_regex,
        "match_all": match_all,
        "use_proximity": use_proximity,
        "proximity": proximity,
        "use_context": use_context,
        "context_before": context_before,
        "context_after": context_after,
    }

    matches = []
    skipped_files = []
    total = len(all_files)
    bar_width = 40
    spinner_chars = "⠋⠙⠹⠸⠼⠴⠦⠧⠇⠏"

    try:
        term_width = os.get_terminal_size().columns
    except OSError:
        term_width = 80

    spinner_lock = threading.Lock()
    spinner_stop = threading.Event()
    spinner_state = {"done": 0, "filename": ""}

    def _render_progress(done, total_count, filename, spinner=""):
        if total_count == 0:
            return
        pct = done / total_count
        filled = int(bar_width * pct)
        bar = "█" * filled + "░" * (bar_width - filled)
        if spinner:
            line = f"\r  [{bar}] {done}/{total_count} {filename} {spinner}"
        else:
            line = f"\r  [{bar}] {done}/{total_count} {filename}"
        line = line.ljust(term_width)[:term_width]
        sys.stdout.write(line)
        sys.stdout.flush()

    def _spinner_thread():
        idx = 0
        while not spinner_stop.is_set():
            with spinner_lock:
                done = spinner_state["done"]
                filename = spinner_state["filename"]
            _render_progress(done, total, filename, spinner_chars[idx % len(spinner_chars)])
            idx += 1
            spinner_stop.wait(0.15)

    spinner = threading.Thread(target=_spinner_thread, daemon=True)
    spinner.start()

    if len(all_files) < 10 or cores == 1:
        for i, filepath in enumerate(all_files):
            filename = os.path.relpath(filepath, cwd)
            with spinner_lock:
                spinner_state["done"] = i
                spinner_state["filename"] = filename
            _render_progress(i, total, filename)
            file_matches, file_skipped = _process_file((filepath, search_config))
            matches.extend(file_matches)
            skipped_files.extend(file_skipped)
    else:
        with multiprocessing.Pool(processes=cores) as pool:
            result_iter = pool.imap(_process_file, [(f, search_config) for f in all_files])
            for i in range(total):
                # Show the file we're about to wait on BEFORE blocking
                filename = os.path.relpath(all_files[i], cwd)
                with spinner_lock:
                    spinner_state["done"] = i
                    spinner_state["filename"] = filename
                _render_progress(i, total, filename)
                # Now block waiting for this file's result
                file_matches, file_skipped = next(result_iter)
                matches.extend(file_matches)
                skipped_files.extend(file_skipped)

    spinner_stop.set()
    spinner.join()

    if total > 0:
        _render_progress(total, total, "done")
        sys.stdout.write("\n")
        sys.stdout.flush()

    for skipped_name, error_msg in skipped_files:
        print(f"Warning: Could not read {skipped_name} ({error_msg})")

    if skipped_files:
        error_log_path = os.path.join(cwd, "docsearch_errors.log")
        with open(error_log_path, "a") as log_f:
            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            for skipped_name, error_msg in skipped_files:
                log_f.write(f"{timestamp}  Could not read {skipped_name} ({error_msg})\n")

    search_elapsed = time.time() - start_time

    output_path = os.path.join(cwd, "docsearch_results.txt")
    if os.path.exists(output_path):
        os.remove(output_path)
    with open(output_path, "w") as f:
        f.write("Program name: docsearch\n")
        f.write("Program Source: https://github.com/exbuf\n")
        f.write("Overview: Searches all supported file types in current directory for search terms.\n")
        f.write("Supported file types:\n")
        f.write(".cfg, .csv, .docx, .epub, .html, .ini, .json, .log, .md, .ods, .odp, .odt, .pdf, .pptx, .rst,\n")
        f.write(".rtf, .sql, .tex, .toml, .tsv, .txt, .xlsx, .xml, .yaml, .yml\n")
        f.write(f"\nReport Generated On ==> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        if use_regex and match_all:
            report_mode = "REGEX+AND"
        elif use_regex:
            report_mode = "REGEX"
        elif match_all:
            report_mode = "ALL"
        else:
            report_mode = "ANY"
        f.write(f"Command ==> {command_str}\n")
        f.write(f"Search Term(s) ==> {', '.join(search_terms)} (match: {report_mode})\n")
        f.write(f"Hits ==> {len(matches)}\n")
        f.write(f"Search Time ==> {search_elapsed:.2f} seconds, Cores used ==> {cores} of {cpu_count}\n")
        total_bytes = sum(os.path.getsize(f_path) for f_path in all_files)
        if total_bytes >= 1_000_000:
            size_str = f"{total_bytes / 1_000_000:.2f} MB"
        elif total_bytes >= 1_000:
            size_str = f"{total_bytes / 1_000:.2f} KB"
        else:
            size_str = f"{total_bytes} bytes"
        f.write(f"Files searched ==> {len(all_files)} ({size_str})\n")
        ext_counts = {}
        for fp in all_files:
            ext = os.path.splitext(fp)[1].lower()
            ext_counts[ext] = ext_counts.get(ext, 0) + 1
        if ext_counts:
            tally = ", ".join(f"{ext}: {count}" for ext, count in ext_counts.items())
            f.write(f"File Types Searched ==> {tally}\n")
        f.write("\n")
        prev_file = None
        for file_dir, filename, line_num, text in matches:
            current_file = os.path.join(file_dir, filename)
            if use_context and prev_file == current_file:
                f.write("---\n\n")
            prev_file = current_file
            if use_context:
                lines = text.split("\n")
                wrapped_lines = [textwrap.fill(line, width=80) if line else line for line in lines]
                wrapped = "\n".join(wrapped_lines)
            else:
                highlighted = text
                for term in search_terms:
                    pattern = term if use_regex else re.escape(term)
                    highlighted = re.sub(pattern, lambda m: f"**{m.group()}**", highlighted, flags=re.IGNORECASE)
                wrapped = textwrap.fill(highlighted, width=80)
            f.write(f'Document: {filename}, Line: {line_num}, Match:\n({file_dir})\n"{wrapped}"\n\n')

    # Create docsearch_results.docx with yellow-highlighted matches
    docx_output_path = os.path.join(cwd, "docsearch_results.docx")
    if os.path.exists(docx_output_path):
        os.remove(docx_output_path)
    result_doc = Document()
    with open(output_path, "r") as f:
        for line in f:
            line = line.rstrip("\n")
            para = result_doc.add_paragraph()

            # Make URL a clickable hyperlink
            if line.startswith("Program Source: "):
                prefix = "Program Source: "
                url = line[len(prefix):]
                para.add_run(prefix)
                r_id = result_doc.part.relate_to(
                    url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True
                )
                hyperlink = OxmlElement("w:hyperlink")
                hyperlink.set(qn("r:id"), r_id)
                run_elem = OxmlElement("w:r")
                rPr = OxmlElement("w:rPr")
                color = OxmlElement("w:color")
                color.set(qn("w:val"), "0000FF")
                rPr.append(color)
                u = OxmlElement("w:u")
                u.set(qn("w:val"), "single")
                rPr.append(u)
                run_elem.append(rPr)
                text_elem = OxmlElement("w:t")
                text_elem.text = url
                run_elem.append(text_elem)
                hyperlink.append(run_elem)
                para._p.append(hyperlink)
                continue

            if line.startswith("Search Term(s) ==> "):
                prefix = "Search Term(s) ==> "
                rest = line[len(prefix):]
                # rest looks like "budget, revenue (match: ANY)"
                match = re.match(r"(.+?)( \(match: \w+\))$", rest)
                if match:
                    terms_str, mode_str = match.group(1), match.group(2)
                    para.add_run(prefix)
                    run = para.add_run(terms_str)
                    run.font.highlight_color = WD_COLOR_INDEX.BRIGHT_GREEN
                    para.add_run(mode_str)
                else:
                    para.add_run(line)
                continue

            is_doc_line = line.startswith("Document:")
            parts = re.split(r"(\*\*.*?\*\*)", line)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    run = para.add_run(part[2:-2])
                    run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                    if is_doc_line:
                        run.bold = True
                else:
                    run = para.add_run(part)
                    if is_doc_line:
                        run.bold = True
    result_doc.save(docx_output_path)

    # Insert report file sizes into both files (after timestamp, before Command)
    def fmt_size(b):
        if b >= 1_000_000:
            return f"{b / 1_000_000:.2f} MB"
        elif b >= 1_000:
            return f"{b / 1_000:.2f} KB"
        return f"{b} bytes"

    txt_size = os.path.getsize(output_path)
    docx_size = os.path.getsize(docx_output_path)
    sizes_line = f"Report File Sizes ==> docsearch_results.txt ({fmt_size(txt_size)}), docsearch_results.docx ({fmt_size(docx_size)})"

    # Update txt report
    with open(output_path, "r") as f:
        content = f.read()
    content = content.replace(
        "\nCommand ==>",
        f"\n{sizes_line}\nCommand ==>",
        1,
    )
    with open(output_path, "w") as f:
        f.write(content)

    # Update docx report — insert sizes paragraph after timestamp
    for i, para in enumerate(result_doc.paragraphs):
        if para.text.startswith("Report Generated On ==>"):
            new_para = OxmlElement("w:p")
            run_elem = OxmlElement("w:r")
            text_elem = OxmlElement("w:t")
            text_elem.text = sizes_line
            run_elem.append(text_elem)
            new_para.append(run_elem)
            para._p.addnext(new_para)
            break
    result_doc.save(docx_output_path)

    if append_name is not None:
        append_txt_path = os.path.join(cwd, f"DO_NOT_SEARCH_ACCUMULATED_{append_name}.txt")
        append_docx_path = os.path.join(cwd, f"DO_NOT_SEARCH_ACCUMULATED_{append_name}.docx")
        with open(output_path, "r") as src:
            results_content = src.read()
        with open(append_txt_path, "a") as dst:
            dst.write(results_content)
        if os.path.exists(append_docx_path):
            existing_doc = Document(append_docx_path)
            new_doc = Document(docx_output_path)
            body = existing_doc.element.body
            sect_pr = body.find(qn('w:sectPr'))
            for para in new_doc.paragraphs:
                new_elem = deepcopy(para._p)
                if sect_pr is not None:
                    sect_pr.addprevious(new_elem)
                else:
                    body.append(new_elem)
            existing_doc.save(append_docx_path)
        else:
            shutil.copy2(docx_output_path, append_docx_path)

    elapsed = time.time() - start_time
    print()
    print(f"Files searched: {len(all_files)} ({size_str})")
    print(f"Found {len(matches)} match(es).")
    print(f"Results ==> {cwd}")
    print(f"  docsearch_results.txt ({fmt_size(txt_size)}), docsearch_results.docx ({fmt_size(docx_size)})")
    if append_name is not None:
        print(f"Results appended to DO_NOT_SEARCH_ACCUMULATED_{append_name}.txt and DO_NOT_SEARCH_ACCUMULATED_{append_name}.docx")
    print(f"Elapsed time: {elapsed:.2f} seconds, Cores used: {cores} of {cpu_count}")
    if skipped_files:
        print(f"Errors logged to docsearch_errors.log ({len(skipped_files)} error(s))")
    print()
    return 0 if matches else 1


if __name__ == "__main__":
    sys.exit(main())
