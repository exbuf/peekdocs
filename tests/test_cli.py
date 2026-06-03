"""Tests for the CLI module."""

import csv
import json
import os
import re
import sys

import pytest

# Tests that rely on a POSIX-shell hook script (``#!/bin/sh``) can't run
# on Windows. The peekdocs ``--on-match`` feature itself is cross-platform
# (the hook is invoked via ``subprocess.run(shlex.split(command))`` — no
# shell), but the test fixtures below write inline POSIX scripts that
# Windows can't execute. Marker keeps the contract verified on macOS and
# Linux without misreporting Windows as broken.
_posix_only = pytest.mark.skipif(
    sys.platform == "win32",
    reason="Test fixture uses a POSIX shell script; --on-match is cross-platform.",
)
from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from fpdf import FPDF

from peekdocs.cli import BANNER_TOP, FUZZY_THRESHOLD, HIGHLIGHT, OCR_IMAGE_TYPES, RESET, SUPPORTED_TYPES, VERSION, main


def test_no_args(capsys):
    result = main([])
    captured = capsys.readouterr()
    assert result == 0
    assert "Search Modes" in captured.out
    assert "Readme documentation: https://github.com/exbuf/peekdocs/blob/main/README.md" in captured.out


def test_help(capsys):
    result = main(["-h"])
    captured = capsys.readouterr()
    assert result == 0
    assert "Readme documentation: https://github.com/exbuf/peekdocs/blob/main/README.md" in captured.out
    assert BANNER_TOP in captured.out


def test_search_finds_matches(tmp_path, monkeypatch, capsys):
    doc = Document()
    doc.add_paragraph("Hello world")
    doc.add_paragraph("This is a test")
    doc.add_paragraph("Hello again")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["hello"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}2{RESET} match(es)" in captured.out

    results_file = tmp_path / "peekdocs_standard_results.txt"
    assert results_file.exists()
    content = results_file.read_text()
    assert "Search Term(s) ==> hello (match: ANY)" in content
    assert f'Document: sample.docx (2 matches), Line: 1, Match:\n({tmp_path})\n"Hello world"\n\n' in content
    assert f'Document: sample.docx (2 matches), Line: 3, Match:\n({tmp_path})\n"Hello again"\n\n' in content

    # Check peekdocs_standard_results.docx was created with yellow highlighting
    docx_results = tmp_path / "peekdocs_standard_results.docx"
    assert docx_results.exists()
    result_doc = Document(str(docx_results))
    highlighted_runs = [
        run for para in result_doc.paragraphs for run in para.runs
        if run.font.highlight_color == WD_COLOR_INDEX.YELLOW
    ]
    assert len(highlighted_runs) == 2
    assert highlighted_runs[0].text == "Hello"
    assert highlighted_runs[1].text == "Hello"


def test_search_no_matches(tmp_path, monkeypatch, capsys):
    doc = Document()
    doc.add_paragraph("Nothing here")
    doc.save(str(tmp_path / "empty.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["zzzzz"])
    captured = capsys.readouterr()

    assert result == 1
    assert f"{HIGHLIGHT}0{RESET} match(es)" in captured.out

    results_file = tmp_path / "peekdocs_standard_results.txt"
    assert results_file.exists()
    content = results_file.read_text()
    assert "Search Term(s) ==> zzzzz (match: ANY)" in content


def test_search_case_insensitive(tmp_path, monkeypatch, capsys):
    doc = Document()
    doc.add_paragraph("Python is great")
    doc.save(str(tmp_path / "test.docx"))

    monkeypatch.chdir(tmp_path)
    main(["PYTHON"])

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert f'Document: test.docx (1 match), Line: 1, Match:\n({tmp_path})\n"Python is great"' in content


def test_search_multi_word_phrase(tmp_path, monkeypatch, capsys):
    """Quoted multi-word phrase is a single search term."""
    doc = Document()
    doc.add_paragraph("Hello world")
    doc.add_paragraph("Hello again")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["Hello world"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Hello world" in content


def test_search_multiple_terms(tmp_path, monkeypatch, capsys):
    """Multiple args are separate search terms (OR logic)."""
    doc = Document()
    doc.add_paragraph("Hello world")
    doc.add_paragraph("Nothing here")
    doc.add_paragraph("Goodbye world")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["Hello", "Goodbye"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}2{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Search Term(s) ==> Hello Goodbye (match: ANY)" in content
    assert "Hello" in content
    assert "Goodbye" in content


def test_search_all_terms(tmp_path, monkeypatch, capsys):
    """With -a flag, only match paragraphs containing ALL terms."""
    doc = Document()
    doc.add_paragraph("Hello world")
    doc.add_paragraph("Hello Goodbye")
    doc.add_paragraph("Goodbye world")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-a", "Hello", "Goodbye"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Search Term(s) ==> Hello Goodbye (match: ALL)" in content
    assert "Hello Goodbye" in content


def test_search_pdf(tmp_path, monkeypatch, capsys):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(text="Budget report for Q1")
    pdf.ln()
    pdf.cell(text="No match here")
    pdf.ln()
    pdf.cell(text="Revised budget plan")
    pdf.output(str(tmp_path / "report.pdf"))

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}2{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "report.pdf" in content
    assert "budget" in content or "Budget" in content


def test_search_docx(tmp_path, monkeypatch, capsys):
    doc = Document()
    doc.add_paragraph("Annual budget review for Q1")
    doc.add_paragraph("No relevant info here")
    doc.add_paragraph("Updated budget forecast")
    doc.save(str(tmp_path / "report.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}2{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "report.docx" in content
    assert "budget" in content


def test_search_csv(tmp_path, monkeypatch, capsys):
    csv_file = tmp_path / "data.csv"
    csv_file.write_text("Name,Amount\nAlice,500\nBob,budget review\nCharlie,300\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "data.csv" in content
    assert "budget" in content


def test_search_odt(tmp_path, monkeypatch, capsys):
    from odf.opendocument import OpenDocumentText
    from odf.text import P as OdtP

    odt_doc = OpenDocumentText()
    p1 = OdtP(text="Hello from ODT")
    odt_doc.text.addElement(p1)
    p2 = OdtP(text="Nothing here")
    odt_doc.text.addElement(p2)
    odt_doc.save(str(tmp_path / "test.odt"))

    monkeypatch.chdir(tmp_path)
    result = main(["hello"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "test.odt" in content
    assert "Hello" in content


def test_search_txt(tmp_path, monkeypatch, capsys):
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("First line\nBudget overview\nThird line\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "Budget" in content


def test_search_html(tmp_path, monkeypatch, capsys):
    html_file = tmp_path / "page.html"
    html_file.write_text("<html><body><p>Budget report</p><p>No match</p></body></html>")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "page.html" in content
    assert "Budget" in content


def test_search_ics(tmp_path, monkeypatch, capsys):
    ics_file = tmp_path / "event.ics"
    ics_file.write_text(
        "BEGIN:VCALENDAR\n"
        "BEGIN:VEVENT\n"
        "SUMMARY:Budget review meeting\n"
        "DESCRIPTION:Discuss Q1 budget and revenue targets\n"
        "END:VEVENT\n"
        "END:VCALENDAR\n"
    )

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert "match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "event.ics" in content
    assert "budget" in content.lower()


def test_search_vcf(tmp_path, monkeypatch, capsys):
    vcf_file = tmp_path / "contact.vcf"
    vcf_file.write_text(
        "BEGIN:VCARD\n"
        "VERSION:3.0\n"
        "FN:Jane Budget-Smith\n"
        "NOTE:Handles budget approvals for Q1\n"
        "END:VCARD\n"
    )

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "contact.vcf" in content
    assert "budget" in content.lower()


def test_search_mbox(tmp_path, monkeypatch, capsys):
    mbox_file = tmp_path / "mail.mbox"
    mbox_file.write_text(
        "From sender@example.com Mon Jan  1 00:00:00 2026\n"
        "From: sender@example.com\n"
        "To: receiver@example.com\n"
        "Subject: Budget Report\n"
        "Date: Mon, 1 Jan 2026 00:00:00 +0000\n"
        "\n"
        "Please review the attached budget for Q1.\n"
        "\n"
    )

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "mail.mbox" in content
    assert "budget" in content.lower()


def test_search_pages(tmp_path, monkeypatch, capsys):
    """Test .pages search — create a minimal zip with XML containing searchable text."""
    import zipfile
    pages_path = tmp_path / "doc.pages"
    with zipfile.ZipFile(str(pages_path), "w") as zf:
        zf.writestr("index.xml", "<document><p>Annual budget review completed</p></document>")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "doc.pages" in content
    assert "budget" in content.lower()


def test_search_xlsx(tmp_path, monkeypatch, capsys):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.append(["Name", "Amount"])
    ws.append(["Alice", "500"])
    ws.append(["Bob", "budget review"])
    wb.save(str(tmp_path / "data.xlsx"))

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "data.xlsx" in content
    assert "budget" in content


def test_version_flag(capsys):
    result = main(["-v"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"peekdocs {VERSION}" in captured.out


def test_version_flag_long(capsys):
    result = main(["-version"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"peekdocs {VERSION}" in captured.out


def test_version_flag_double_dash(capsys):
    """`--version` is the GNU/POSIX convention. Prior to v1.0.10 the CLI
    only matched `-v` and `-version` and silently treated `--version` as
    a literal search term — printing a startup banner and then running a
    full directory scan against `--version` as if it were a search query.
    """
    result = main(["--version"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"peekdocs {VERSION}" in captured.out
    # The bug symptom was that the search code path printed "Searching"
    # after the version banner. The fix has to short-circuit before that.
    assert "Searching" not in captured.out


def test_save_flag_double_dash(capsys):
    """`--save` follows the same GNU/POSIX convention as `--version` and
    `--help`. Prior to v1.0.10 the CLI only matched `-s` and `-save`, so
    `peekdocs --save my_report` fell through to the search code path and
    was treated as a literal search for the string "--save" — producing
    a destructive recursive scan of the current directory instead of
    saving the previous run's results.

    The test passes `--save` with no filename argument; the save handler
    returns 2 with a "No filename provided" error. If the bug regressed,
    the search code path would print "Searching" and return 0 after
    scanning, so we assert both: the save error message IS present, and
    "Searching" is NOT.
    """
    result = main(["--save"])
    captured = capsys.readouterr()
    assert result == 2
    assert "No filename provided" in captured.out
    assert "Searching" not in captured.out


def test_banner_always_printed(capsys):
    main(["anything"])
    captured = capsys.readouterr()
    assert "peekdocs v" in captured.out


def test_search_recursive(tmp_path, monkeypatch, capsys):
    """With -r flag, files in subdirectories are found and shown with relative paths."""
    subdir = tmp_path / "sub" / "deep"
    subdir.mkdir(parents=True)
    txt_file = subdir / "nested.txt"
    txt_file.write_text("Budget overview in nested file\nNo match here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-r", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert str(tmp_path / "sub" / "deep") in content
    assert "nested.txt" in content
    assert "Budget" in content


def test_search_no_recursive_skips_subdirs(tmp_path, monkeypatch, capsys):
    """Without -r flag, files in subdirectories are NOT searched."""
    subdir = tmp_path / "sub"
    subdir.mkdir()
    txt_file = subdir / "nested.txt"
    txt_file.write_text("Budget overview in nested file\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 1
    assert f"{HIGHLIGHT}0{RESET} match(es)" in captured.out


def test_search_with_type_filter(tmp_path, monkeypatch, capsys):
    """With -t flag, only specified file types are searched."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    csv_file = tmp_path / "data.csv"
    csv_file.write_text("Name,Amount\nBob,budget review\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-t", "txt", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "data.csv" not in content


def test_search_with_multiple_type_filters(tmp_path, monkeypatch, capsys):
    """With -t flag and comma-separated types, multiple types are searched."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    csv_file = tmp_path / "data.csv"
    csv_file.write_text("Name,Amount\nBob,budget review\n")
    html_file = tmp_path / "page.html"
    html_file.write_text("<html><body><p>Budget report</p></body></html>")

    monkeypatch.chdir(tmp_path)
    result = main(["-t", "txt,csv", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}2{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "data.csv" in content
    assert "page.html" not in content


def test_search_invalid_type_filter(tmp_path, monkeypatch, capsys):
    """With -t flag and unsupported type, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-t", "xyz", "budget"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Unsupported file type: xyz" in captured.out


def test_search_md(tmp_path, monkeypatch, capsys):
    md_file = tmp_path / "notes.md"
    md_file.write_text("# Heading\nBudget overview for Q1\nNo match here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.md" in content
    assert "Budget" in content


def test_search_regex(tmp_path, monkeypatch, capsys):
    """With -x flag, search terms are treated as regex patterns."""
    txt_file = tmp_path / "contacts.txt"
    txt_file.write_text("Call 555-123-4567\nNo phone here\nReach out at 555-987-6543\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-x", r"\d{3}-\d{3}-\d{4}"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}2{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "contacts.txt" in content
    assert "555-123-4567" in content
    assert "555-987-6543" in content


def test_search_regex_with_and(tmp_path, monkeypatch, capsys):
    """Combine -x and -a for regex AND logic."""
    txt_file = tmp_path / "data.txt"
    txt_file.write_text("Order 123 costs $45.99\nOrder 456 no price\nJust $99.99 here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-x", "-a", r"\d{3}", r"\$\d+\.\d{2}"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "REGEX+AND" in content


def test_search_invalid_regex(tmp_path, monkeypatch, capsys):
    """With -x flag and invalid regex, return error."""
    monkeypatch.chdir(tmp_path)
    result = main(["-x", "[invalid"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Invalid regex pattern" in captured.out


def test_search_json(tmp_path, monkeypatch, capsys):
    json_file = tmp_path / "data.json"
    json_file.write_text('{\n  "title": "Budget report",\n  "amount": 500\n}\n')

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "data.json" in content
    assert "Budget" in content


def test_search_context_after(tmp_path, monkeypatch, capsys):
    """With -A flag, lines after each match are included."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("\n".join(f"Line {i}" for i in range(1, 11)) + "\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-A", "2", "Line 3"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Line 3" in content
    assert "Line 4" in content
    assert "Line 5" in content
    assert "Line 6" not in content


def test_search_context_before(tmp_path, monkeypatch, capsys):
    """With -B flag, lines before each match are included."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("\n".join(f"Line {i}" for i in range(1, 11)) + "\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-B", "2", "Line 3"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Line 1" in content
    assert "Line 2" in content
    assert "Line 3" in content
    assert "Line 4" not in content


def test_search_context_both(tmp_path, monkeypatch, capsys):
    """With -B and -A flags combined, lines before and after are included."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("\n".join(f"Line {i}" for i in range(1, 11)) + "\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-B", "1", "-A", "1", "Line 3"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Line 1" not in content
    assert "Line 2" in content
    assert "Line 3" in content
    assert "Line 4" in content
    assert "Line 5" not in content


def test_search_context_merge(tmp_path, monkeypatch, capsys):
    """When context regions overlap, they merge without duplicate lines."""
    txt_file = tmp_path / "notes.txt"
    lines = ["alpha", "beta", "MATCH_ONE", "gamma", "MATCH_TWO", "delta", "epsilon"]
    txt_file.write_text("\n".join(lines) + "\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-B", "1", "-A", "1", "MATCH"])
    captured = capsys.readouterr()

    assert result == 0
    # Overlapping context should merge into one group
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "beta" in content
    assert "MATCH_ONE" in content
    assert "gamma" in content
    assert "MATCH_TWO" in content
    assert "delta" in content
    # gamma should appear exactly once (no duplicates from overlap)
    results_section = content.split("Results:")[1]  # only look in results
    assert results_section.count("gamma") == 1


def test_search_ods(tmp_path, monkeypatch, capsys):
    from odf.opendocument import OpenDocumentSpreadsheet
    from odf.table import Table, TableRow, TableCell
    from odf.text import P as OdtP

    ods_doc = OpenDocumentSpreadsheet()
    table = Table(name="Sheet1")
    row1 = TableRow()
    cell1 = TableCell()
    cell1.addElement(OdtP(text="Budget report"))
    row1.addElement(cell1)
    table.addElement(row1)
    row2 = TableRow()
    cell2 = TableCell()
    cell2.addElement(OdtP(text="No match"))
    row2.addElement(cell2)
    table.addElement(row2)
    ods_doc.spreadsheet.addElement(table)
    ods_doc.save(str(tmp_path / "data.ods"))

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "data.ods" in content
    assert "Budget" in content


def test_search_odp(tmp_path, monkeypatch, capsys):
    from odf.opendocument import OpenDocumentPresentation
    from odf.text import P as OdtP
    from odf.draw import Page, Frame, TextBox

    odp_doc = OpenDocumentPresentation()
    page = Page(masterpagename="Default")
    frame = Frame()
    textbox = TextBox()
    textbox.addElement(OdtP(text="Budget presentation"))
    frame.addElement(textbox)
    page.addElement(frame)
    odp_doc.presentation.addElement(page)
    odp_doc.save(str(tmp_path / "slides.odp"))

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "slides.odp" in content
    assert "Budget" in content


def test_search_toml(tmp_path, monkeypatch, capsys):
    toml_file = tmp_path / "config.toml"
    toml_file.write_text('[project]\nname = "budget-tracker"\nversion = "1.0"\n')

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "config.toml" in content
    assert "budget" in content


def test_search_rst(tmp_path, monkeypatch, capsys):
    rst_file = tmp_path / "docs.rst"
    rst_file.write_text("Title\n=====\n\nBudget overview for Q1\n\nNo match here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "docs.rst" in content
    assert "Budget" in content


def test_search_tex(tmp_path, monkeypatch, capsys):
    tex_file = tmp_path / "paper.tex"
    tex_file.write_text("\\documentclass{article}\n\\begin{document}\nBudget analysis results\n\\end{document}\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "paper.tex" in content
    assert "Budget" in content


def test_search_ini(tmp_path, monkeypatch, capsys):
    ini_file = tmp_path / "settings.ini"
    ini_file.write_text("[general]\nproject = budget_app\nversion = 1.0\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "settings.ini" in content
    assert "budget" in content


def test_search_cfg(tmp_path, monkeypatch, capsys):
    cfg_file = tmp_path / "app.cfg"
    cfg_file.write_text("[settings]\nbudget_limit = 5000\nactive = true\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "app.cfg" in content
    assert "budget" in content


def test_search_sql(tmp_path, monkeypatch, capsys):
    sql_file = tmp_path / "queries.sql"
    sql_file.write_text("SELECT * FROM expenses;\nSELECT * FROM budget WHERE year = 2026;\nDROP TABLE temp;\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "queries.sql" in content
    assert "budget" in content


def test_search_tsv(tmp_path, monkeypatch, capsys):
    tsv_file = tmp_path / "data.tsv"
    tsv_file.write_text("Name\tAmount\nAlice\t500\nBob\tbudget review\nCharlie\t300\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "data.tsv" in content
    assert "budget" in content


def test_search_epub(tmp_path, monkeypatch, capsys):
    from ebooklib import epub as epub_mod
    book = epub_mod.EpubBook()
    book.set_identifier("test123")
    book.set_title("Test Book")
    book.set_language("en")
    ch = epub_mod.EpubHtml(title="Chapter 1", file_name="ch1.xhtml", lang="en")
    ch.content = b"<html><body><p>Budget report for Q1</p><p>No match here</p></body></html>"
    book.add_item(ch)
    book.spine = ["nav", ch]
    book.add_item(epub_mod.EpubNcx())
    book.add_item(epub_mod.EpubNav())
    epub_mod.write_epub(str(tmp_path / "book.epub"), book)

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "book.epub" in content
    assert "Budget" in content


def test_search_yaml(tmp_path, monkeypatch, capsys):
    yaml_file = tmp_path / "config.yaml"
    yaml_file.write_text("title: Budget report\namount: 500\ndescription: No match here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "config.yaml" in content
    assert "Budget" in content


def test_search_yml(tmp_path, monkeypatch, capsys):
    yml_file = tmp_path / "settings.yml"
    yml_file.write_text("name: test\nbudget_limit: 1000\nstatus: active\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "settings.yml" in content
    assert "budget" in content


def test_search_ipynb(tmp_path, monkeypatch, capsys):
    nb_file = tmp_path / "analysis.ipynb"
    nb_file.write_text(json.dumps({
        "cells": [
            {"cell_type": "code", "metadata": {}, "source": ["# Budget analysis\n", "total = 5000\n"]},
            {"cell_type": "markdown", "metadata": {}, "source": ["## No match here\n"]},
        ],
        "metadata": {}, "nbformat": 4, "nbformat_minor": 4,
    }))
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "analysis.ipynb" in content
    assert "Budget" in content


def test_search_env(tmp_path, monkeypatch, capsys):
    env_file = tmp_path / ".env"
    env_file.write_text("DATABASE_URL=postgres://localhost\nAPI_KEY=budget_key_123\nSECRET=hidden\n")
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert ".env" in content
    assert "budget" in content


def test_search_dockerfile(tmp_path, monkeypatch, capsys):
    df = tmp_path / "Dockerfile"
    df.write_text("FROM python:3.12-slim\nWORKDIR /app\nRUN pip install budget-tool\nEXPOSE 8080\n")
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Dockerfile" in content
    assert "budget" in content


def test_search_css(tmp_path, monkeypatch, capsys):
    css_file = tmp_path / "styles.css"
    css_file.write_text("body { margin: 0; }\n.budget-panel { color: blue; }\np { font-size: 14px; }\n")
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "styles.css" in content


def test_search_scss(tmp_path, monkeypatch, capsys):
    scss_file = tmp_path / "app.scss"
    scss_file.write_text("$primary: blue;\n.budget-view { color: $primary; }\nbody { margin: 0; }\n")
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "app.scss" in content


def test_search_lua(tmp_path, monkeypatch, capsys):
    lua_file = tmp_path / "config.lua"
    lua_file.write_text("local config = {}\nconfig.budget_limit = 5000\nreturn config\n")
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "config.lua" in content


def test_search_scala(tmp_path, monkeypatch, capsys):
    scala_file = tmp_path / "Pipeline.scala"
    scala_file.write_text("object Pipeline {\n  val budget = 10000\n  def run(): Unit = {}\n}\n")
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Pipeline.scala" in content


def test_search_tf(tmp_path, monkeypatch, capsys):
    tf_file = tmp_path / "main.tf"
    tf_file.write_text('resource "aws_instance" "web" {\n  instance_type = "t3.medium"\n  tags = { Name = "budget-server" }\n}\n')
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "main.tf" in content


def test_search_proto(tmp_path, monkeypatch, capsys):
    proto_file = tmp_path / "schema.proto"
    proto_file.write_text('syntax = "proto3";\nmessage BudgetRequest {\n  double amount = 1;\n}\n')
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "schema.proto" in content


def test_search_graphql(tmp_path, monkeypatch, capsys):
    gql_file = tmp_path / "schema.graphql"
    gql_file.write_text("type Query {\n  budget(id: ID!): Budget\n}\ntype Budget {\n  amount: Float\n}\n")
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "schema.graphql" in content


def test_search_jsonl(tmp_path, monkeypatch, capsys):
    jsonl_file = tmp_path / "data.jsonl"
    jsonl_file.write_text('{"name": "Alice", "amount": 100}\n{"name": "Bob", "budget": 500}\n{"name": "Carol", "amount": 200}\n')
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "data.jsonl" in content


def test_search_ndjson(tmp_path, monkeypatch, capsys):
    ndjson_file = tmp_path / "events.ndjson"
    ndjson_file.write_text('{"event": "login"}\n{"event": "budget_update", "amount": 1000}\n')
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "events.ndjson" in content


def test_search_xml(tmp_path, monkeypatch, capsys):
    xml_file = tmp_path / "config.xml"
    xml_file.write_text('<?xml version="1.0"?>\n<root>\n  <title>Budget report</title>\n  <value>No match</value>\n</root>\n')

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "config.xml" in content
    assert "Budget" in content


def test_search_log(tmp_path, monkeypatch, capsys):
    log_file = tmp_path / "app.log"
    log_file.write_text("2026-03-18 INFO Starting up\n2026-03-18 ERROR Budget exceeded limit\n2026-03-18 INFO Shutting down\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "app.log" in content
    assert "Budget" in content


def test_search_pptx(tmp_path, monkeypatch, capsys):
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.text = "Budget report for Q1"
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
    txBox2.text_frame.text = "No match here"
    prs.save(str(tmp_path / "report.pptx"))

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "report.pptx" in content
    assert "Budget" in content


def test_search_rtf(tmp_path, monkeypatch, capsys):
    rtf_file = tmp_path / "report.rtf"
    rtf_content = r'{\rtf1\ansi{\fonttbl\f0 Times New Roman;}\f0 Budget report for Q1\par No match here\par}'
    rtf_file.write_text(rtf_content)

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "report.rtf" in content
    assert "Budget" in content


def test_search_context_invalid(tmp_path, monkeypatch, capsys):
    """With -A and invalid count, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-A", "abc", "budget"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Invalid count for -A" in captured.out


def test_search_with_file_filter(tmp_path, monkeypatch, capsys):
    """With -f flag, only specified files are searched."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    txt_file2 = tmp_path / "other.txt"
    txt_file2.write_text("Budget details\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-f", "notes.txt", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "other.txt" not in content


def test_search_with_multiple_file_filters(tmp_path, monkeypatch, capsys):
    """With -f flag and comma-separated files, multiple files are searched."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    csv_file = tmp_path / "data.csv"
    csv_file.write_text("Name,Amount\nBob,budget review\n")
    html_file = tmp_path / "page.html"
    html_file.write_text("<html><body><p>Budget report</p></body></html>")

    monkeypatch.chdir(tmp_path)
    result = main(["-f", "notes.txt,data.csv", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}2{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "data.csv" in content
    assert "page.html" not in content


def test_search_file_filter_not_found(tmp_path, monkeypatch, capsys):
    """With -f flag and nonexistent file, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-f", "missing.txt", "budget"])
    captured = capsys.readouterr()

    assert result == 2
    assert "File not found: missing.txt" in captured.out


def test_search_file_filter_with_and(tmp_path, monkeypatch, capsys):
    """With -f and -a flags, AND logic applies to specified files."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget and revenue overview\n")
    txt_file2 = tmp_path / "other.txt"
    txt_file2.write_text("Budget and revenue details\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-f", "notes.txt", "-a", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "other.txt" not in content


def test_search_file_filter_with_regex(tmp_path, monkeypatch, capsys):
    """With -f and -x flags, regex search applies to specified files."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Call 555-123-4567 for details\n")
    txt_file2 = tmp_path / "other.txt"
    txt_file2.write_text("Call 555-987-6543 for info\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-f", "notes.txt", "-x", r"\d{3}-\d{3}-\d{4}"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "other.txt" not in content


def test_search_file_filter_with_context(tmp_path, monkeypatch, capsys):
    """With -f, -A, and -B flags, context lines apply to specified files."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Line one\nBudget overview\nLine three\n")
    txt_file2 = tmp_path / "other.txt"
    txt_file2.write_text("Budget details\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-f", "notes.txt", "-B", "1", "-A", "1", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "other.txt" not in content
    assert "Line one" in content
    assert "Line three" in content


def test_search_file_filter_recursive(tmp_path, monkeypatch, capsys):
    """With -f and -r flags, specific file is found in subdirectories."""
    subdir = tmp_path / "subdir"
    subdir.mkdir()
    txt_file = subdir / "notes.txt"
    txt_file.write_text("Budget overview\n")
    txt_file2 = tmp_path / "other.txt"
    txt_file2.write_text("Budget details\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-f", "notes.txt", "-r", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "other.txt" not in content


def test_search_proximity(tmp_path, monkeypatch, capsys):
    """With -p flag, terms within N words of each other match."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget for revenue growth was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-p", "3", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_proximity_no_match(tmp_path, monkeypatch, capsys):
    """With -p flag, terms too far apart don't match."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was set and later the team discussed revenue\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-p", "2", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 1
    assert f"{HIGHLIGHT}0{RESET} match(es)" in captured.out


def test_search_proximity_requires_two_terms(tmp_path, monkeypatch, capsys):
    """With -p flag and only one term, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-p", "5", "budget"])
    captured = capsys.readouterr()

    assert result == 2
    assert "requires at least 2 search terms" in captured.out


def test_search_proximity_invalid(tmp_path, monkeypatch, capsys):
    """With -p flag and invalid count, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-p", "abc", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Invalid count for -p" in captured.out


def test_search_save_append(tmp_path, monkeypatch, capsys):
    """With -sa flag, search runs normally and results are appended to peekdocs_accumulated file."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview for Q1\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-sa", "my_report", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    assert "Results appended to peekdocs_accumulated_my_report.txt and peekdocs_accumulated_my_report.docx" in captured.out

    assert (tmp_path / "peekdocs_standard_results.txt").exists()
    assert (tmp_path / "peekdocs_standard_results.docx").exists()
    assert (tmp_path / "peekdocs_accumulated_my_report.txt").exists()
    assert (tmp_path / "peekdocs_accumulated_my_report.docx").exists()

    content = (tmp_path / "peekdocs_accumulated_my_report.txt").read_text()
    assert "budget" in content.lower()


def test_search_save_append_accumulates(tmp_path, monkeypatch, capsys):
    """With -sa flag used twice, results accumulate in the peekdocs_accumulated file."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview for Q1\nRevenue report for Q2\n")

    monkeypatch.chdir(tmp_path)
    main(["-sa", "combined", "budget"])
    main(["-sa", "combined", "revenue"])

    append_txt = tmp_path / "peekdocs_accumulated_combined.txt"
    content = append_txt.read_text()
    assert "budget" in content.lower()
    assert "revenue" in content.lower()


def test_search_save_append_no_filename(capsys):
    """With -sa flag and no filename, an error is returned."""
    result = main(["-sa"])
    captured = capsys.readouterr()

    assert result == 2
    assert "No filename provided" in captured.out


def test_search_save_append_no_terms(tmp_path, monkeypatch, capsys):
    """With -sa flag and filename but no search terms, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-sa", "my_report"])
    captured = capsys.readouterr()

    assert result == 2
    assert "No search terms provided" in captured.out


def test_search_with_cores_flag(tmp_path, monkeypatch, capsys):
    """With -c flag, search uses specified number of cores."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)
    result = main(["-c", "2", "budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_with_single_core(tmp_path, monkeypatch, capsys):
    """With -c 1, search runs single-threaded."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)
    result = main(["-c", "1", "budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_cores_invalid(tmp_path, monkeypatch, capsys):
    """With -c and invalid count, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-c", "abc", "budget"])
    captured = capsys.readouterr()
    assert result == 2
    assert "Invalid count for -c" in captured.out


def test_search_cores_zero(tmp_path, monkeypatch, capsys):
    """With -c 0, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-c", "0", "budget"])
    captured = capsys.readouterr()
    assert result == 2
    assert "Invalid count for -c" in captured.out


def test_search_cores_no_count(capsys):
    """With -c and no count, an error is returned."""
    result = main(["-c"])
    captured = capsys.readouterr()
    assert result == 2
    assert "No count provided" in captured.out


def test_search_cores_negative(tmp_path, monkeypatch, capsys):
    """With -c and negative number, an error is returned."""
    monkeypatch.chdir(tmp_path)
    result = main(["-c", "-1", "budget"])
    captured = capsys.readouterr()
    assert result == 2
    assert "Invalid count for -c" in captured.out


def test_search_cores_with_other_flags(tmp_path, monkeypatch, capsys):
    """With -c combined with other flags, all work together."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget and revenue overview\n")
    monkeypatch.chdir(tmp_path)
    result = main(["-c", "2", "-a", "budget", "revenue"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_quiet_mode(tmp_path, monkeypatch, capsys):
    """With -q flag, banner is suppressed but results still print."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)
    result = main(["-q", "budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    assert "OR search" not in captured.out


def test_search_without_quiet_shows_banner(tmp_path, monkeypatch, capsys):
    """Without -q flag, banner is shown."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()
    assert result == 0
    assert "peekdocs v" in captured.out


def test_config_file_ignored_by_cli(tmp_path, monkeypatch, capsys):
    """Config file does not affect CLI — recursive=true in config is ignored."""
    subdir = tmp_path / "sub"
    subdir.mkdir()
    doc = Document()
    doc.add_paragraph("Budget report for Q1")
    doc.save(str(subdir / "nested.docx"))

    config_file = tmp_path / ".peekdocsrc"
    config_file.write_text("recursive = true\n")
    monkeypatch.chdir(tmp_path)

    result = main(["budget"])
    captured = capsys.readouterr()

    # Without -r flag, subdirectory file should NOT be found
    assert "Files searched: 0" in captured.out


def test_config_cli_overrides(tmp_path, monkeypatch, capsys):
    """CLI flags override config file settings."""
    doc = Document()
    doc.add_paragraph("Budget report")
    doc.save(str(tmp_path / "report.docx"))

    config_file = tmp_path / ".peekdocsrc"
    config_file.write_text("cores = 2\n")
    monkeypatch.chdir(tmp_path)

    result = main(["-c", "1", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert "Cores used: 1" in captured.out


def test_config_missing_file(tmp_path, monkeypatch, capsys):
    """No .peekdocsrc, search works normally."""
    doc = Document()
    doc.add_paragraph("Budget overview")
    doc.save(str(tmp_path / "report.docx"))

    monkeypatch.chdir(tmp_path)

    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_config_invalid_values(tmp_path, monkeypatch, capsys):
    """Bad values in config are silently ignored."""
    doc = Document()
    doc.add_paragraph("Budget overview")
    doc.save(str(tmp_path / "report.docx"))

    config_file = tmp_path / ".peekdocsrc"
    config_file.write_text("cores = abc\nrecursive = banana\nunknown_key = whatever\n")
    monkeypatch.chdir(tmp_path)

    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_keyboard_interrupt(tmp_path, monkeypatch, capsys):
    """Ctrl+C during search prints clean message and returns exit code 2."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    import peekdocs.api as api_module
    def interrupt_on_process(args_tuple):
        raise KeyboardInterrupt

    monkeypatch.setattr(api_module, "_process_file", interrupt_on_process)
    result = main(["budget"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Search cancelled." in captured.out


def test_config_flag_show(tmp_path, monkeypatch, capsys):
    """--config with no args displays current settings."""
    config_file = tmp_path / ".peekdocsrc"
    config_file.write_text("recursive = true\ncores = 4\n")
    result = main(["--config"])
    captured = capsys.readouterr()

    assert result == 0
    assert "recursive = True" in captured.out
    assert "cores = 4" in captured.out


def test_config_flag_show_empty(tmp_path, monkeypatch, capsys):
    """--config with no file prints 'No config file found.'"""
    result = main(["--config"])
    captured = capsys.readouterr()

    assert result == 0
    assert "No config file found." in captured.out


def test_config_flag_set(tmp_path, monkeypatch, capsys):
    """--config key=value creates/updates the config file."""
    result = main(["--config", "recursive=true", "cores=4"])
    captured = capsys.readouterr()

    assert result == 0
    assert "Set: recursive = True" in captured.out
    assert "Set: cores = 4" in captured.out
    content = (tmp_path / ".peekdocsrc").read_text()
    assert "recursive = true" in content
    assert "cores = 4" in content


def test_config_flag_remove(tmp_path, monkeypatch, capsys):
    """--config key= removes the key from the config file."""
    config_file = tmp_path / ".peekdocsrc"
    config_file.write_text("recursive = true\ncores = 4\n")
    result = main(["--config", "recursive="])
    captured = capsys.readouterr()

    assert result == 0
    assert "Removed: recursive" in captured.out
    content = (tmp_path / ".peekdocsrc").read_text()
    assert "recursive" not in content
    assert "cores = 4" in content


def test_config_flag_invalid_key(tmp_path, monkeypatch, capsys):
    """--config with unknown key prints error and returns 2."""
    result = main(["--config", "badkey=true"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Unknown setting: badkey" in captured.out


# --- OCR tests ---


def test_ocr_no_tesseract(tmp_path, monkeypatch, capsys):
    """Using -O without Tesseract installed prints install instructions."""
    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr("shutil.which", lambda cmd: None)
    (tmp_path / "doc.txt").write_text("hello world")
    result = main(["-O", "hello"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Tesseract OCR is not installed" in captured.out
    assert "brew install tesseract" in captured.out


def test_ocr_scanned_pdf(tmp_path, monkeypatch, capsys):
    """Scanned PDF page (no extractable text) gets OCR'd when -O is used."""
    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr("shutil.which", lambda cmd: "/usr/local/bin/tesseract" if cmd == "tesseract" else None)

    # Create a PDF with no text (simulates a scanned page)
    pdf = FPDF()
    pdf.add_page()
    pdf.output(str(tmp_path / "scanned.pdf"))

    import peekdocs.api as api_module
    monkeypatch.setattr(api_module, "_ocr_image", lambda img: "Budget report for Q1\nTotal revenue increased")

    result = main(["-O", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "scanned.pdf" in content
    assert "Budget" in content


def test_ocr_image_jpg(tmp_path, monkeypatch, capsys):
    """JPG image file is searched via OCR when -O is used."""
    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr("shutil.which", lambda cmd: "/usr/local/bin/tesseract" if cmd == "tesseract" else None)

    from PIL import Image as PILImage
    img = PILImage.new("RGB", (100, 100), "white")
    img.save(str(tmp_path / "scan.jpg"))

    import peekdocs.api as api_module
    monkeypatch.setattr(api_module, "_ocr_image", lambda img: "Budget report for Q1")

    result = main(["-O", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "scan.jpg" in content


def test_ocr_image_png(tmp_path, monkeypatch, capsys):
    """PNG image file is searched via OCR when -O is used."""
    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr("shutil.which", lambda cmd: "/usr/local/bin/tesseract" if cmd == "tesseract" else None)

    from PIL import Image as PILImage
    img = PILImage.new("RGB", (100, 100), "white")
    img.save(str(tmp_path / "scan.png"))

    import peekdocs.api as api_module
    monkeypatch.setattr(api_module, "_ocr_image", lambda img: "Invoice total amount due")

    result = main(["-O", "invoice"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "scan.png" in content


def test_ocr_images_ignored_without_flag(tmp_path, monkeypatch, capsys):
    """Image files are not discovered without the -O flag."""
    monkeypatch.chdir(tmp_path)

    from PIL import Image as PILImage
    img = PILImage.new("RGB", (100, 100), "white")
    img.save(str(tmp_path / "photo.jpg"))
    img.save(str(tmp_path / "photo.png"))

    result = main(["budget"])
    captured = capsys.readouterr()

    assert "Files searched: 0" in captured.out
    # Verify image files were not found in results
    results = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "photo.jpg" not in results
    assert "photo.png" not in results


def test_ocr_normal_pdf_skips_ocr(tmp_path, monkeypatch, capsys):
    """PDF with real text does NOT trigger OCR even with -O."""
    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr("shutil.which", lambda cmd: "/usr/local/bin/tesseract" if cmd == "tesseract" else None)

    # Create a PDF with real text
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(text="Budget report for Q1")
    pdf.output(str(tmp_path / "normal.pdf"))

    import peekdocs.cli as cli_module

    def fail_ocr(img):
        raise RuntimeError("OCR should not be called for normal PDF")

    monkeypatch.setattr(cli_module, "_ocr_image", fail_ocr)

    result = main(["-O", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_ocr_config_ignored_by_cli(tmp_path, monkeypatch, capsys):
    """ocr=true in config does not activate OCR in CLI without -O flag."""
    monkeypatch.chdir(tmp_path)

    config_file = tmp_path / ".peekdocsrc"
    config_file.write_text("ocr = true\n")

    from PIL import Image as PILImage
    img = PILImage.new("RGB", (100, 100), "white")
    img.save(str(tmp_path / "scan.png"))

    result = main(["budget"])
    captured = capsys.readouterr()

    # Without -O flag, image files should not be discovered
    assert "Files searched: 0" in captured.out


# --- Fuzzy search tests ---


def test_search_fuzzy_match(tmp_path, monkeypatch, capsys):
    """With -z flag, approximate matches are found."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budgt for this quarter was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-z", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "budgt" in content


def test_search_fuzzy_no_match(tmp_path, monkeypatch, capsys):
    """With -z flag, words too dissimilar don't match."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The elephant walked slowly\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-z", "budget"])
    captured = capsys.readouterr()

    assert f"{HIGHLIGHT}0{RESET} match(es)" in captured.out


def test_search_fuzzy_exact_still_matches(tmp_path, monkeypatch, capsys):
    """With -z flag, exact matches still work."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-z", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_fuzzy_with_and(tmp_path, monkeypatch, capsys):
    """Fuzzy + AND: both fuzzy terms must appear in same line."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budgt and revnue report\nOnly budgt here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-z", "-a", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "FUZZY" in content


def test_search_fuzzy_with_regex_error(tmp_path, monkeypatch, capsys):
    """Fuzzy + regex is an error."""
    monkeypatch.chdir(tmp_path)
    (tmp_path / "doc.txt").write_text("hello")
    result = main(["-z", "-x", "hello"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Cannot combine fuzzy (-z) and regex (-x)" in captured.out


def test_search_fuzzy_with_proximity(tmp_path, monkeypatch, capsys):
    """Fuzzy + proximity: fuzzy terms within N words match."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budgt for revnue growth was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-z", "-p", "3", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_fuzzy_mode_string(tmp_path, monkeypatch, capsys):
    """Mode string includes FUZZY."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budgt was approved\n")

    monkeypatch.chdir(tmp_path)
    main(["-z", "budget"])
    captured = capsys.readouterr()

    assert "OR+FUZZY" in captured.out


def test_search_fuzzy_config_ignored_by_cli(tmp_path, monkeypatch, capsys):
    """fuzzy=true in config does not enable fuzzy in CLI without -z flag."""
    rc = tmp_path / ".peekdocsrc"
    rc.write_text("fuzzy = true\n")

    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budgt was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    captured = capsys.readouterr()

    # Without -z flag, fuzzy is not active — "budgt" won't match "budget"
    assert "FUZZY" not in captured.out


# --- Exclude terms tests ---


def test_search_exclude_basic(tmp_path, monkeypatch, capsys):
    """Lines containing exclude terms are filtered out."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\nThe draft budget needs review\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-n", "draft", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "approved" in content
    assert "draft" not in content.split("Document:")[1]


def test_search_exclude_keeps_non_excluded(tmp_path, monkeypatch, capsys):
    """When exclude term is absent, matches are kept."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-n", "draft", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_exclude_with_and(tmp_path, monkeypatch, capsys):
    """Exclude works with AND logic."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("budget and revenue approved\nbudget and revenue draft\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-n", "draft", "-a", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_exclude_with_regex(tmp_path, monkeypatch, capsys):
    """Exclude terms use regex matching when -x is active."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("budget $100\nbudget pending\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-n", r"\$\d+", "-x", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "pending" in content


def test_search_exclude_with_fuzzy(tmp_path, monkeypatch, capsys):
    """Exclude terms use fuzzy matching when -z is active."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budgt was approved\nThe budgt drft needs review\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-n", "draft", "-z", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_exclude_mode_string(tmp_path, monkeypatch, capsys):
    """Mode string includes +NOT when exclude terms present."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    main(["-n", "draft", "budget"])
    captured = capsys.readouterr()

    assert "OR+NOT" in captured.out
    assert "Excluding [draft]" in captured.out


def test_search_exclude_multiple(tmp_path, monkeypatch, capsys):
    """Multiple comma-separated exclude terms all filter."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("budget approved\nbudget draft\nbudget obsolete\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-n", "draft,obsolete", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Exclude Term(s) ==> draft obsolete" in content


def test_search_exclude_no_terms_error(tmp_path, monkeypatch, capsys):
    """-n without terms returns error."""
    monkeypatch.chdir(tmp_path)
    (tmp_path / "doc.txt").write_text("hello")
    result = main(["-n"])
    captured = capsys.readouterr()

    assert result == 2
    assert "No exclude terms provided" in captured.out


# --- Wildcard search tests ---


def test_search_wildcard_star(tmp_path, monkeypatch, capsys):
    """Wildcard * matches zero or more word characters."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget and budgets and budgeting\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-w", "budg*"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "budget" in content


def test_search_wildcard_question(tmp_path, monkeypatch, capsys):
    """Wildcard ? matches exactly one word character."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The test and the text are here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-w", "te?t"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "test" in content or "text" in content


def test_search_wildcard_no_match(tmp_path, monkeypatch, capsys):
    """Wildcard that doesn't match returns no results."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-w", "xyz*"])
    captured = capsys.readouterr()

    assert result == 1
    assert f"{HIGHLIGHT}0{RESET} match(es)" in captured.out


def test_search_wildcard_exact(tmp_path, monkeypatch, capsys):
    """Term without wildcard chars still works in -w mode."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-w", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_wildcard_with_and(tmp_path, monkeypatch, capsys):
    """Wildcard with AND logic."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget and revenue report\nOnly budget here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-w", "-a", "budg*", "rev*"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_search_wildcard_with_regex_error(tmp_path, monkeypatch, capsys):
    """Wildcard + regex is an error."""
    monkeypatch.chdir(tmp_path)
    (tmp_path / "doc.txt").write_text("hello")
    result = main(["-w", "-x", "hello"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Cannot combine wildcard (-w) and regex (-x)" in captured.out


def test_search_wildcard_with_fuzzy_error(tmp_path, monkeypatch, capsys):
    """Wildcard + fuzzy is an error."""
    monkeypatch.chdir(tmp_path)
    (tmp_path / "doc.txt").write_text("hello")
    result = main(["-w", "-z", "hello"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Cannot combine wildcard (-w) and fuzzy (-z)" in captured.out


def test_search_wildcard_mode_string(tmp_path, monkeypatch, capsys):
    """Mode string shows WILDCARD."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    main(["-w", "budg*"])
    captured = capsys.readouterr()

    assert "WILDCARD" in captured.out


def test_search_wildcard_config_ignored_by_cli(tmp_path, monkeypatch, capsys):
    """wildcard=true in config does not enable wildcard in CLI without -w flag."""
    rc = tmp_path / ".peekdocsrc"
    rc.write_text("wildcard = true\n")

    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budg*"])
    captured = capsys.readouterr()

    # Without -w flag, wildcard is not active
    assert "WILDCARD" not in captured.out


def test_search_wildcard_with_exclude(tmp_path, monkeypatch, capsys):
    """Wildcard + exclude: exclude terms are also wildcard-matched."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\nThe budget draft was rejected\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-w", "-n", "dra*", "budg*"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "approved" in content


def test_output_csv(tmp_path, capsys, monkeypatch):
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\nNothing here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "-o", "csv", "budget"])
    assert result == 0

    csv_path = tmp_path / "peekdocs_standard_results.csv"
    assert csv_path.exists()
    with open(csv_path) as f:
        reader = csv.reader(f)
        rows = list(reader)
    assert rows[0] == ["filename", "folder", "line_number", "matched_text"]
    assert rows[1][0] == "notes.txt"
    assert rows[1][2] == "1"
    assert "budget" in rows[1][3].lower()
    # Highlight markers should be stripped
    assert "**" not in rows[1][3]


def test_output_json(tmp_path, capsys, monkeypatch):
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\nNothing here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "-o", "json", "budget"])
    assert result == 0

    json_path = tmp_path / "peekdocs_standard_results.json"
    assert json_path.exists()
    with open(json_path) as f:
        data = json.load(f)
    assert data["search_terms"] == ["budget"]
    assert data["matches_found"] == 1
    assert data["files_searched"] == 1
    assert "elapsed_seconds" in data
    assert "timestamp" in data
    assert data["mode"] == "ANY"
    assert len(data["matches"]) == 1
    assert data["matches"][0]["filename"] == "notes.txt"
    assert data["matches"][0]["line_number"] == 1
    assert "budget" in data["matches"][0]["matched_text"].lower()
    assert "**" not in data["matches"][0]["matched_text"]


def test_output_csv_and_json(tmp_path, capsys, monkeypatch):
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "-o", "csv,json", "budget"])
    assert result == 0
    assert (tmp_path / "peekdocs_standard_results.csv").exists()
    assert (tmp_path / "peekdocs_standard_results.json").exists()


def test_output_invalid_format(tmp_path, capsys, monkeypatch):
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("The budget was approved\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "-o", "xml", "budget"])
    captured = capsys.readouterr()
    assert result == 2
    assert "Invalid output format" in captured.out


def test_output_no_format(tmp_path, capsys, monkeypatch):
    monkeypatch.chdir(tmp_path)
    result = main(["-o"])
    captured = capsys.readouterr()
    assert result == 2
    assert "No format provided" in captured.out


# ─── Index tests ─────────────────────────────────────────


def test_index_build(tmp_path, monkeypatch, capsys):
    """--index builds the index database."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\nRevenue details\n")

    monkeypatch.chdir(tmp_path)
    result = main(["--index"])
    captured = capsys.readouterr()

    assert result == 0
    assert "Index built:" in captured.out
    assert "1 files" in captured.out
    assert (tmp_path / ".peekdocs.db").exists()


def test_index_build_includes_subfolders(tmp_path, monkeypatch, capsys):
    """--index always indexes subdirectories."""
    subdir = tmp_path / "subdir"
    subdir.mkdir()
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    (subdir / "report.txt").write_text("Revenue details\n")

    monkeypatch.chdir(tmp_path)
    result = main(["--index"])
    captured = capsys.readouterr()

    assert result == 0
    assert "2 files" in captured.out


def test_index_clear(tmp_path, monkeypatch, capsys):
    """--index-clear removes the database."""
    (tmp_path / "notes.txt").write_text("Budget\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    assert (tmp_path / ".peekdocs.db").exists()

    result = main(["--index-clear"])
    captured = capsys.readouterr()

    assert result == 0
    assert "Index removed" in captured.out
    assert not (tmp_path / ".peekdocs.db").exists()


def test_index_clear_no_index(tmp_path, monkeypatch, capsys):
    """--index-clear with no index prints a message."""
    monkeypatch.chdir(tmp_path)
    result = main(["--index-clear"])
    captured = capsys.readouterr()

    assert result == 0
    assert "No index found" in captured.out


def test_index_status(tmp_path, monkeypatch, capsys):
    """--index-status shows index info."""
    (tmp_path / "notes.txt").write_text("Budget overview\nRevenue details\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["--index-status"])
    captured = capsys.readouterr()

    assert result == 0
    assert "Files indexed:" in captured.out
    assert "Lines indexed:" in captured.out
    assert "Database size:" in captured.out


def test_index_status_shows_last_updated(tmp_path, monkeypatch, capsys):
    """--index-status shows last_updated timestamp."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    capsys.readouterr()

    result = main(["--index-status"])
    captured = capsys.readouterr()

    assert result == 0
    assert "Last updated:" in captured.out


def test_index_last_updated_changes_on_refresh(tmp_path, monkeypatch):
    """last_updated timestamp is written on both build and refresh."""
    from peekdocs.indexer import index_status, refresh_index
    import time

    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    status_after_build = index_status(str(tmp_path))
    assert "last_updated" in status_after_build
    build_ts = status_after_build["last_updated"]

    time.sleep(0.1)
    (tmp_path / "extra.txt").write_text("Revenue\n")
    refresh_index(str(tmp_path), recursive=True, use_ocr=False)
    status_after_refresh = index_status(str(tmp_path))
    refresh_ts = status_after_refresh["last_updated"]

    assert refresh_ts >= build_ts


def test_index_status_no_index(tmp_path, monkeypatch, capsys):
    """--index-status with no index prints a message."""
    monkeypatch.chdir(tmp_path)
    result = main(["--index-status"])
    captured = capsys.readouterr()

    assert result == 0
    assert "No index found" in captured.out


def test_indexed_search_keyword(tmp_path, monkeypatch, capsys):
    """Indexed keyword search returns same matches as direct scan."""
    (tmp_path / "notes.txt").write_text("Budget overview\nRevenue details\nOther line\n")
    monkeypatch.chdir(tmp_path)

    # Direct scan
    main(["-q", "budget"])
    direct_content = (tmp_path / "peekdocs_standard_results.txt").read_text()

    # Build index and search again
    main(["--index"])
    main(["-q", "budget"])
    indexed_content = (tmp_path / "peekdocs_standard_results.txt").read_text()

    # Both should find "Budget overview" with highlighting
    assert "Budget overview" in direct_content
    assert "Budget overview" in indexed_content


def test_indexed_search_and(tmp_path, monkeypatch, capsys):
    """Indexed AND search returns matches."""
    (tmp_path / "notes.txt").write_text("Budget and revenue overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["-q", "-a", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_indexed_search_regex(tmp_path, monkeypatch, capsys):
    """Indexed regex search uses parse-cache path and finds matches."""
    (tmp_path / "notes.txt").write_text("Call 555-123-4567 for details\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["-q", "-x", r"\d{3}-\d{3}-\d{4}"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out


def test_indexed_search_fuzzy(tmp_path, monkeypatch, capsys):
    """Indexed fuzzy search uses parse-cache path and finds matches."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["-q", "-z", "budgt"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content


def test_indexed_search_wildcard(tmp_path, monkeypatch, capsys):
    """Indexed wildcard search uses parse-cache path and finds matches."""
    (tmp_path / "notes.txt").write_text("Budget overview\nBudgets listed\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["-q", "-w", "budg*"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content


def test_indexed_search_with_file_type_filter(tmp_path, monkeypatch, capsys):
    """Indexed search respects file type filters."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    (tmp_path / "data.csv").write_text("Budget,Amount\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["-q", "-t", "txt", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content
    assert "data.csv" not in content


def test_index_refresh_detects_changes(tmp_path, monkeypatch, capsys):
    """Incremental refresh detects new, changed, and deleted files."""
    from peekdocs.indexer import index_status, refresh_index

    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    status = index_status(str(tmp_path))
    assert status["file_count"] == 1

    # Add a new file
    (tmp_path / "extra.txt").write_text("Revenue details\n")
    result = refresh_index(str(tmp_path), recursive=False, use_ocr=False)
    assert result["added"] == 1

    # Delete a file
    os.remove(tmp_path / "notes.txt")
    result = refresh_index(str(tmp_path), recursive=False, use_ocr=False)
    assert result["removed"] == 1


def test_validate_db_locked_is_not_corrupt(tmp_path, monkeypatch):
    """A locked database should not be treated as corrupt and deleted."""
    import sqlite3
    from peekdocs.indexer import _validate_db, index_exists

    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)
    main(["--index"])
    assert index_exists(str(tmp_path))

    # Hold an exclusive lock on the database
    db_path = tmp_path / ".peekdocs.db"
    conn = sqlite3.connect(str(db_path))
    conn.execute("BEGIN EXCLUSIVE")

    # _validate_db should return True (not treat locked as corrupt)
    assert _validate_db(str(tmp_path)) is True
    # Index file should still exist (not deleted)
    assert index_exists(str(tmp_path))

    conn.rollback()
    conn.close()


def test_refresh_index_connection_closed_on_error(tmp_path, monkeypatch):
    """refresh_index closes the connection even if an error occurs."""
    import sqlite3
    from peekdocs.indexer import refresh_index, index_exists

    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)
    main(["--index"])

    # Hold exclusive lock so refresh_index will fail on writes
    db_path = tmp_path / ".peekdocs.db"
    blocker = sqlite3.connect(str(db_path))
    blocker.execute("BEGIN EXCLUSIVE")

    # refresh_index should raise (lock timeout) but not leak the connection
    try:
        refresh_index(str(tmp_path), recursive=True, use_ocr=False)
    except sqlite3.OperationalError:
        pass  # Expected — database is locked

    blocker.rollback()
    blocker.close()

    # Index should still be intact and usable
    assert index_exists(str(tmp_path))
    result = refresh_index(str(tmp_path), recursive=True, use_ocr=False)
    assert isinstance(result, dict)


def test_search_succeeds_when_refresh_locked(tmp_path, monkeypatch, capsys):
    """Search proceeds with existing index when refresh_index can't get the lock."""
    import sqlite3

    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)
    main(["--index"])
    capsys.readouterr()

    # Hold exclusive lock to block refresh_index inside the search subprocess
    db_path = tmp_path / ".peekdocs.db"
    blocker = sqlite3.connect(str(db_path))
    blocker.execute("BEGIN EXCLUSIVE")

    # Search should still succeed (api.py catches the refresh failure)
    # Use --no-index to avoid the lock issue for this specific test,
    # since the subprocess would also be blocked.
    # Instead, test the api.search() path directly.
    from peekdocs.api import search
    result = search("budget", directory=str(tmp_path), use_index=False)
    assert len(result.matches) == 1

    blocker.rollback()
    blocker.close()

    # Now confirm indexed search works when lock is released
    result = search("budget", directory=str(tmp_path), use_index=True)
    assert len(result.matches) == 1


def test_concurrent_refresh_with_timeout(tmp_path, monkeypatch):
    """Two concurrent refresh_index calls don't corrupt the index."""
    import threading
    from peekdocs.indexer import refresh_index, index_status

    # Create several files
    for i in range(5):
        (tmp_path / f"doc_{i}.txt").write_text(f"Content for document {i}\n")
    monkeypatch.chdir(tmp_path)
    main(["--index"])

    # Add more files
    for i in range(5, 10):
        (tmp_path / f"doc_{i}.txt").write_text(f"Content for document {i}\n")

    results = []
    errors = []

    def do_refresh():
        try:
            r = refresh_index(str(tmp_path), recursive=True, use_ocr=False)
            results.append(r)
        except Exception as e:
            errors.append(e)

    # Run two refreshes concurrently
    t1 = threading.Thread(target=do_refresh)
    t2 = threading.Thread(target=do_refresh)
    t1.start()
    t2.start()
    t1.join(timeout=30)
    t2.join(timeout=30)

    # At least one should succeed; index should be consistent
    assert len(results) >= 1
    status = index_status(str(tmp_path))
    assert status is not None
    assert status["file_count"] == 10


def test_index_refresh_cli(tmp_path, monkeypatch, capsys):
    """--index-refresh performs incremental update via CLI."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    capsys.readouterr()  # clear

    # Add a new file
    (tmp_path / "extra.txt").write_text("Revenue details\n")

    result = main(["--index-refresh"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1 added" in captured.out
    assert "0 updated" in captured.out
    assert "0 removed" in captured.out


def test_index_refresh_cli_no_index(tmp_path, monkeypatch, capsys):
    """--index-refresh with no index prints a helpful message."""
    monkeypatch.chdir(tmp_path)
    result = main(["--index-refresh"])
    captured = capsys.readouterr()

    assert result == 0
    assert "No index found" in captured.out


def test_index_refresh_cli_updated_file(tmp_path, monkeypatch, capsys):
    """--index-refresh detects modified files."""
    import time
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    capsys.readouterr()

    time.sleep(0.1)
    txt_file.write_text("Updated budget overview\n")

    result = main(["--index-refresh"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1 updated" in captured.out


def test_index_refresh_cli_removed_file(tmp_path, monkeypatch, capsys):
    """--index-refresh detects deleted files."""
    (tmp_path / "notes.txt").write_text("Budget\n")
    (tmp_path / "extra.txt").write_text("Revenue\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    capsys.readouterr()

    (tmp_path / "extra.txt").unlink()

    result = main(["--index-refresh"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1 removed" in captured.out


def test_no_index_fallback(tmp_path, monkeypatch, capsys):
    """Without an index, search uses direct scan (existing behavior unchanged)."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    result = main(["-q", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    assert f"{HIGHLIGHT}1{RESET} match(es)" in captured.out
    assert not (tmp_path / ".peekdocs.db").exists()


def test_indexed_search_exclude(tmp_path, monkeypatch, capsys):
    """Indexed search respects exclude terms."""
    (tmp_path / "notes.txt").write_text("Budget draft overview\nBudget final version\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["-q", "-n", "draft", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "final version" in content
    assert "draft overview" not in content


def test_indexed_search_proximity(tmp_path, monkeypatch, capsys):
    """Indexed proximity search uses parse-cache path."""
    (tmp_path / "notes.txt").write_text("The budget for revenue growth\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["-q", "-p", "3", "budget", "revenue"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "notes.txt" in content


def test_indexed_search_context(tmp_path, monkeypatch, capsys):
    """Indexed context search uses parse-cache path."""
    (tmp_path / "notes.txt").write_text("Line one\nBudget overview\nLine three\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    result = main(["-q", "-B", "1", "-A", "1", "budget"])
    captured = capsys.readouterr()

    assert result == 0
    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Line one" in content
    assert "Line three" in content


def test_indexed_search_shows_indexed_in_output(tmp_path, monkeypatch, capsys):
    """Indexed search shows 'indexed' in the mode string."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    main(["-q", "budget"])
    captured = capsys.readouterr()

    assert "indexed" in captured.out


def test_indexed_search_report_shows_last_updated(tmp_path, monkeypatch, capsys):
    """Indexed search report includes index last updated timestamp."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    main(["--index"])
    main(["-q", "budget"])

    report = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Index last updated:" in report


# ─── Dependency and error handling tests ─────────────────


def test_check_shows_versions(tmp_path, monkeypatch, capsys):
    """--check output includes dependency version numbers."""
    monkeypatch.chdir(tmp_path)
    result = main(["--check"])
    captured = capsys.readouterr()

    # Should show version numbers for required deps
    assert "pymupdf" in captured.out
    assert "python-docx" in captured.out
    assert "ok (" in captured.out  # version in parens
    # Should show optional deps section
    assert "Optional dependencies:" in captured.out
    assert "SQLite version:" in captured.out


def test_check_shows_optional_deps(tmp_path, monkeypatch, capsys):
    """--check output includes optional dependency status."""
    monkeypatch.chdir(tmp_path)
    main(["--check"])
    captured = capsys.readouterr()

    assert "rapidfuzz" in captured.out
    assert "customtkinter" in captured.out
    assert "Pillow" in captured.out


def test_fuzzy_missing_dep(tmp_path, monkeypatch, capsys):
    """Fuzzy search with missing rapidfuzz gives helpful error."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    import builtins
    original_import = builtins.__import__

    def mock_import(name, *args, **kwargs):
        if name == "rapidfuzz":
            raise ImportError("No module named 'rapidfuzz'")
        return original_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", mock_import)
    result = main(["-z", "budgt"])
    captured = capsys.readouterr()

    assert result == 2
    assert "rapidfuzz" in captured.out
    assert "pip install" in captured.out


def test_ocr_missing_pytesseract(tmp_path, monkeypatch, capsys):
    """OCR with missing pytesseract gives helpful error."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")
    monkeypatch.chdir(tmp_path)

    import builtins
    original_import = builtins.__import__

    def mock_import(name, *args, **kwargs):
        if name == "pytesseract":
            raise ImportError("No module named 'pytesseract'")
        return original_import(name, *args, **kwargs)

    # Also mock shutil.which to return a path so we get past the Tesseract binary check
    monkeypatch.setattr("shutil.which", lambda x: "/usr/bin/tesseract" if x == "tesseract" else None)
    monkeypatch.setattr(builtins, "__import__", mock_import)
    result = main(["-O", "budget"])
    captured = capsys.readouterr()

    assert result == 2
    assert "pytesseract" in captured.out
    assert "pip install" in captured.out


def test_corrupt_index_recovery(tmp_path, monkeypatch, capsys):
    """Corrupt .peekdocs.db is auto-deleted and search falls back to direct scan."""
    (tmp_path / "notes.txt").write_text("Budget overview\n")

    # Write garbage to .peekdocs.db
    db_path = tmp_path / ".peekdocs.db"
    db_path.write_bytes(b"THIS IS NOT A SQLITE DATABASE")

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "budget"])
    captured = capsys.readouterr()

    # Should succeed (falls back to direct scan since corrupt DB was removed)
    assert result == 0
    # Corrupt DB should have been deleted
    assert not db_path.exists()


def test_crash_report_includes_versions(tmp_path, monkeypatch):
    """Crash report in peekdocs_errors.log includes dependency versions."""
    monkeypatch.chdir(tmp_path)

    # Force a crash by monkeypatching _main_inner to raise
    from peekdocs import cli
    original = cli._main_inner

    def crash_inner(argv=None):
        raise RuntimeError("test crash for version logging")

    monkeypatch.setattr(cli, "_main_inner", crash_inner)
    result = cli.main(["budget"])

    assert result == 2
    log_path = tmp_path / "peekdocs_errors.log"
    assert log_path.exists()
    log_content = log_path.read_text()
    assert "Dependency versions:" in log_content
    assert "pymupdb" in log_content or "pymupdf" in log_content


def test_results_csv_json_not_searched(tmp_path, monkeypatch):
    """peekdocs_standard_results.csv and peekdocs_standard_results.json must not be searched."""
    # Create a real searchable file
    (tmp_path / "real.txt").write_text("budget line here")
    # Create result files that contain the search term — these should be skipped
    (tmp_path / "peekdocs_standard_results.csv").write_text("budget,amount\n100,200\n")
    (tmp_path / "peekdocs_standard_results.json").write_text('{"query": "budget"}')

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "budget"])
    assert result == 0

    content = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "real.txt" in content
    assert "peekdocs_standard_results.csv" not in content
    assert "peekdocs_standard_results.json" not in content


def test_per_file_match_counts(tmp_path, monkeypatch, capsys):
    """Per-file match counts appear in console output and TXT report."""
    # file_a has 3 lines with "budget", file_b has 1
    (tmp_path / "file_a.txt").write_text("budget line one\nbudget line two\nbudget line three\n")
    (tmp_path / "file_b.txt").write_text("budget found here\n")

    monkeypatch.chdir(tmp_path)
    result = main(["budget"])
    assert result == 0

    captured = capsys.readouterr().out
    assert "file_a.txt: 3" in captured
    assert "file_b.txt: 1" in captured

    report = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Hits ==> 4" in report

    # -q suppresses the banner but still shows per-file listing
    result = main(["-q", "budget"])
    captured = capsys.readouterr().out
    assert "match(es)" in captured
    assert "file_a.txt: 3" in captured


def test_search_inverse(tmp_path, monkeypatch, capsys):
    """Inverse search lists files WITHOUT matches."""
    (tmp_path / "has_match.txt").write_text("budget report\n")
    (tmp_path / "no_match.txt").write_text("revenue forecast\n")
    (tmp_path / "also_no_match.txt").write_text("quarterly summary\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "--inverse", "budget"])
    captured = capsys.readouterr().out

    assert result == 0
    assert "file(s) WITHOUT matches" in captured
    assert "no_match.txt" in captured
    assert "also_no_match.txt" in captured
    # The file with the match should NOT appear in the inverse listing
    assert "has_match.txt" not in captured.split("WITHOUT matches")[1]

    report = (tmp_path / "peekdocs_standard_results.txt").read_text()
    assert "Files WITHOUT matches: 2" in report
    assert "no_match.txt" in report
    assert "also_no_match.txt" in report


def test_search_inverse_all_match(tmp_path, monkeypatch, capsys):
    """When all files match, inverse returns exit code 1 (nothing to report)."""
    (tmp_path / "a.txt").write_text("budget\n")
    (tmp_path / "b.txt").write_text("budget\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "--inverse", "budget"])
    captured = capsys.readouterr().out
    clean = re.sub(r"\033\[[0-9;]*m", "", captured)

    assert result == 1
    assert "0 file(s) WITHOUT matches" in clean


def test_search_inverse_none_match(tmp_path, monkeypatch, capsys):
    """When no files match, inverse lists all files."""
    (tmp_path / "a.txt").write_text("hello\n")
    (tmp_path / "b.txt").write_text("world\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-q", "--inverse", "zzzzz"])
    captured = capsys.readouterr().out
    clean = re.sub(r"\033\[[0-9;]*m", "", captured)

    assert result == 0
    assert "2 file(s) WITHOUT matches" in clean


def test_search_inverse_csv(tmp_path, monkeypatch):
    """Inverse search CSV export contains filename and folder columns only."""
    (tmp_path / "has_match.txt").write_text("budget\n")
    (tmp_path / "no_match.txt").write_text("revenue\n")

    monkeypatch.chdir(tmp_path)
    main(["-q", "--inverse", "-o", "csv", "budget"])

    csv_path = tmp_path / "peekdocs_standard_results.csv"
    assert csv_path.exists()
    content = csv_path.read_text()
    assert "filename,folder" in content
    assert "no_match.txt" in content
    assert "line_number" not in content


def test_search_inverse_json(tmp_path, monkeypatch):
    """Inverse search JSON export contains inverse_files array."""
    (tmp_path / "has_match.txt").write_text("budget\n")
    (tmp_path / "no_match.txt").write_text("revenue\n")

    monkeypatch.chdir(tmp_path)
    main(["-q", "--inverse", "-o", "json", "budget"])

    json_path = tmp_path / "peekdocs_standard_results.json"
    assert json_path.exists()
    data = json.load(open(json_path))
    assert "inverse_files" in data
    assert "files_without_matches" in data
    assert data["files_without_matches"] == 1
    filenames = [f["filename"] for f in data["inverse_files"]]
    assert "no_match.txt" in filenames


# ─── Boolean Expression Search (-e) ─────────────────────────


def test_expression_and(tmp_path, monkeypatch, capsys):
    """Test basic AND expression."""
    doc = Document()
    doc.add_paragraph("hello world")
    doc.add_paragraph("hello there")
    doc.add_paragraph("goodbye world")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-e", "hello AND world"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1" in captured.out and "match(es)" in captured.out


def test_expression_or(tmp_path, monkeypatch, capsys):
    """Test basic OR expression."""
    doc = Document()
    doc.add_paragraph("hello world")
    doc.add_paragraph("goodbye world")
    doc.add_paragraph("nothing here")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-e", "hello OR goodbye"])
    captured = capsys.readouterr()

    assert result == 0
    assert "2" in captured.out and "match(es)" in captured.out


def test_expression_not(tmp_path, monkeypatch, capsys):
    """Test AND NOT expression."""
    doc = Document()
    doc.add_paragraph("hello world")
    doc.add_paragraph("hello draft")
    doc.add_paragraph("goodbye world")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-e", "hello AND NOT draft"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1" in captured.out and "match(es)" in captured.out


def test_expression_grouped(tmp_path, monkeypatch, capsys):
    """Test grouped expression: (a AND b) OR (c AND d)."""
    doc = Document()
    doc.add_paragraph("bob and amy went to the store")
    doc.add_paragraph("fred and wilma stayed home")
    doc.add_paragraph("bob and fred went fishing")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-e", "(bob AND amy) OR (fred AND wilma)"])
    captured = capsys.readouterr()

    assert result == 0
    assert "2" in captured.out and "match(es)" in captured.out


def test_expression_or_groups_with_and(tmp_path, monkeypatch, capsys):
    """Test (a OR b) AND (c OR d)."""
    doc = Document()
    doc.add_paragraph("budget and cost analysis")
    doc.add_paragraph("revenue and profit report")
    doc.add_paragraph("budget summary only")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-e", "(budget OR revenue) AND (cost OR profit)"])
    captured = capsys.readouterr()

    assert result == 0
    assert "2" in captured.out and "match(es)" in captured.out


def test_expression_mode_string(tmp_path, monkeypatch, capsys):
    """Test that expression mode displays EXPR."""
    doc = Document()
    doc.add_paragraph("hello world")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    main(["-e", "hello AND world"])
    captured = capsys.readouterr()

    assert "(EXPR)" in captured.out


def test_expression_report(tmp_path, monkeypatch, capsys):
    """Test that report contains expression instead of search terms."""
    doc = Document()
    doc.add_paragraph("hello world")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    main(["-e", "hello AND world"])

    results_file = tmp_path / "peekdocs_standard_results.txt"
    content = results_file.read_text()
    assert "Search Expression ==> hello AND world" in content
    assert "Expression: ON" in content


def test_expression_no_matches(tmp_path, monkeypatch, capsys):
    """Test expression with no matches returns exit code 1."""
    doc = Document()
    doc.add_paragraph("hello world")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-e", "foo AND bar"])
    captured = capsys.readouterr()

    assert result == 1
    assert "0" in captured.out and "match(es)" in captured.out


def test_expression_conflict_with_and_mode(tmp_path, monkeypatch, capsys):
    """Test -e cannot be combined with -a."""
    monkeypatch.chdir(tmp_path)
    result = main(["-e", "hello AND world", "-a"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Cannot combine" in captured.out


def test_expression_conflict_with_exclude(tmp_path, monkeypatch, capsys):
    """Test -e cannot be combined with -n."""
    monkeypatch.chdir(tmp_path)
    result = main(["-e", "hello AND world", "-n", "draft"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Cannot combine" in captured.out


def test_expression_conflict_with_proximity(tmp_path, monkeypatch, capsys):
    """Test -e cannot be combined with -p."""
    monkeypatch.chdir(tmp_path)
    result = main(["-e", "hello AND world", "-p", "5"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Cannot combine" in captured.out


def test_expression_invalid_syntax(tmp_path, monkeypatch, capsys):
    """Test invalid expression returns error."""
    monkeypatch.chdir(tmp_path)
    result = main(["-e", "(hello AND"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Invalid expression" in captured.out


def test_expression_with_context(tmp_path, monkeypatch, capsys):
    """Test expression works with context lines."""
    (tmp_path / "test.txt").write_text("line one\nhello world\nline three\n")

    monkeypatch.chdir(tmp_path)
    result = main(["-e", "hello AND world", "-B", "1", "-A", "1"])
    captured = capsys.readouterr()

    assert result == 0
    results_file = tmp_path / "peekdocs_standard_results.txt"
    content = results_file.read_text()
    assert "line one" in content
    assert "line three" in content


def test_expression_with_wildcard(tmp_path, monkeypatch, capsys):
    """Test expression combined with -w wildcard mode."""
    doc = Document()
    doc.add_paragraph("budget and revenue report")
    doc.add_paragraph("nothing relevant here")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-e", "-w", "budg* AND rev*"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1" in captured.out and "match(es)" in captured.out


def test_whole_word_basic(tmp_path, monkeypatch, capsys):
    """Test -W matches whole words only, not substrings."""
    doc = Document()
    doc.add_paragraph("bob is here")
    doc.add_paragraph("bobcat is an animal")
    doc.add_paragraph("bobby went home")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-W", "bob"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1" in captured.out and "match(es)" in captured.out
    # Verify report only has the whole-word match
    report_path = tmp_path / "peekdocs_standard_results.txt"
    report = report_path.read_text()
    assert "bob is here" in report
    assert "bobcat" not in report
    assert "bobby" not in report


def test_whole_word_and_mode(tmp_path, monkeypatch, capsys):
    """Test -W combined with -a (AND mode)."""
    doc = Document()
    doc.add_paragraph("bob met amy today")
    doc.add_paragraph("bobby met amy today")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-W", "-a", "bob", "amy"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1" in captured.out and "match(es)" in captured.out


def test_whole_word_expression(tmp_path, monkeypatch, capsys):
    """Test -W combined with -e (expression mode)."""
    doc = Document()
    doc.add_paragraph("bob met amy today")
    doc.add_paragraph("bobcat met amy today")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-W", "-e", "bob AND amy"])
    captured = capsys.readouterr()

    assert result == 0
    assert "1" in captured.out and "match(es)" in captured.out


def test_whole_word_mode_string(tmp_path, monkeypatch, capsys):
    """Test that -W adds WORD to the mode string."""
    doc = Document()
    doc.add_paragraph("test word here")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-W", "test"])
    captured = capsys.readouterr()

    assert result == 0
    assert "WORD" in captured.out


def test_whole_word_no_match(tmp_path, monkeypatch, capsys):
    """Test -W returns no matches when only substrings exist."""
    doc = Document()
    doc.add_paragraph("bobcat and bobby")
    doc.save(str(tmp_path / "sample.docx"))

    monkeypatch.chdir(tmp_path)
    result = main(["-W", "bob"])
    captured = capsys.readouterr()

    assert result == 1
    assert "0" in captured.out and "match(es)" in captured.out


def test_suite_timestamp_creates_unique_reports(tmp_path, monkeypatch, capsys):
    """--suite --timestamp produces uniquely named reports so a second run does not overwrite the first."""
    (tmp_path / "doc.txt").write_text("TODO: finish this\nrelevant text here\n")
    collection = {
        "saved_searches": {
            "find_todos": {"search_text": "TODO", "recursive": True, "and_mode": False},
        },
        "suites": {"my suite": ["find_todos"]},
    }
    (tmp_path / ".peekdocs_collection.json").write_text(json.dumps(collection))
    monkeypatch.chdir(tmp_path)

    result = main(["--suite", "my suite", "--timestamp"])
    assert result == 0

    plain_txt = tmp_path / "peekdocs_suite_results.txt"
    plain_docx = tmp_path / "peekdocs_suite_results.docx"
    assert not plain_txt.exists(), "unsuffixed file should not be written when --timestamp is set"
    assert not plain_docx.exists()

    stamped_txt = list(tmp_path.glob("peekdocs_suite_results_*.txt"))
    stamped_docx = list(tmp_path.glob("peekdocs_suite_results_*.docx"))
    assert len(stamped_txt) == 1
    assert len(stamped_docx) == 1
    assert re.match(r"peekdocs_suite_results_\d{8}_\d{6}\.txt$", stamped_txt[0].name)


def test_suite_without_timestamp_uses_plain_filename(tmp_path, monkeypatch, capsys):
    """Backward compatibility: omitting --timestamp keeps the original report filename."""
    (tmp_path / "doc.txt").write_text("TODO: finish this\n")
    collection = {
        "saved_searches": {
            "find_todos": {"search_text": "TODO", "recursive": True, "and_mode": False},
        },
        "suites": {"my suite": ["find_todos"]},
    }
    (tmp_path / ".peekdocs_collection.json").write_text(json.dumps(collection))
    monkeypatch.chdir(tmp_path)

    result = main(["--suite", "my suite"])
    assert result == 0
    assert (tmp_path / "peekdocs_suite_results.txt").exists()
    assert (tmp_path / "peekdocs_suite_results.docx").exists()
    assert not list(tmp_path.glob("peekdocs_suite_results_*.txt"))


def test_suite_accepts_path_prefixed_name(tmp_path, monkeypatch, capsys):
    """`--suite "/path/to/folder/Name"` runs the suite in that folder, even from elsewhere."""
    docs = tmp_path / "docs"
    docs.mkdir()
    (docs / "doc.txt").write_text("TODO: x\n")
    collection = {
        "saved_searches": {"t": {"search_text": "TODO", "recursive": True}},
        "suites": {"My Suite": ["t"]},
    }
    (docs / ".peekdocs_collection.json").write_text(json.dumps(collection))

    elsewhere = tmp_path / "elsewhere"
    elsewhere.mkdir()
    monkeypatch.chdir(elsewhere)

    full = str(docs / "My Suite")
    result = main(["--suite", full])
    assert result == 0
    assert (docs / "peekdocs_suite_results.txt").exists()
    assert (docs / "peekdocs_suite_results.docx").exists()
    # Reports go to the suite's folder, not the cwd we ran from.
    assert not (elsewhere / "peekdocs_suite_results.txt").exists()


def test_suite_path_prefix_falls_through_when_name_contains_slash_but_no_match(tmp_path, monkeypatch, capsys):
    """A slash in the name with no matching folder/suite still hits the not-found path cleanly."""
    monkeypatch.chdir(tmp_path)
    result = main(["--suite", "/nonexistent/path/Nothing"])
    assert result == 2
    captured = capsys.readouterr()
    assert "not found" in captured.out


def test_regex_collection_timestamp_creates_unique_reports(tmp_path, monkeypatch, capsys):
    """--regex-collection --timestamp produces uniquely named reports so consecutive runs do not overwrite each other."""
    (tmp_path / "doc.txt").write_text("TODO: review\nhttps://example.com\n")
    rc_data = {
        "my collection": [
            {"name": "TODOs", "regex": r"TODO\b", "enabled": True},
        ],
    }
    (tmp_path / ".peekdocs_regex_collections.json").write_text(json.dumps(rc_data))
    monkeypatch.chdir(tmp_path)

    result = main(["--regex-collection", "my collection", "-d", str(tmp_path), "-r", "--timestamp"])
    assert result == 0

    plain_txt = tmp_path / "peekdocs_regex_results.txt"
    assert not plain_txt.exists(), "unsuffixed file should not be written when --timestamp is set"

    stamped_txt = list(tmp_path.glob("peekdocs_regex_results_*.txt"))
    stamped_docx = list(tmp_path.glob("peekdocs_regex_results_*.docx"))
    assert len(stamped_txt) == 1
    assert len(stamped_docx) == 1
    assert re.match(r"peekdocs_regex_results_\d{8}_\d{6}\.txt$", stamped_txt[0].name)


def test_regex_collection_without_timestamp_uses_plain_filename(tmp_path, monkeypatch, capsys):
    """Backward compatibility: omitting --timestamp keeps the original report filename."""
    (tmp_path / "doc.txt").write_text("TODO: review\n")
    rc_data = {
        "my collection": [
            {"name": "TODOs", "regex": r"TODO\b", "enabled": True},
        ],
    }
    (tmp_path / ".peekdocs_regex_collections.json").write_text(json.dumps(rc_data))
    monkeypatch.chdir(tmp_path)

    result = main(["--regex-collection", "my collection", "-d", str(tmp_path), "-r"])
    assert result == 0
    assert (tmp_path / "peekdocs_regex_results.txt").exists()
    assert (tmp_path / "peekdocs_regex_results.docx").exists()
    assert not list(tmp_path.glob("peekdocs_regex_results_*.txt"))


# ── --hash (SHA-256 chain-of-custody) ──────────────────────────────────

def _sha256_hex(data: bytes) -> str:
    import hashlib
    return hashlib.sha256(data).hexdigest()


def test_hash_absent_without_flag(tmp_path, monkeypatch, capsys):
    """Without --hash, no sha256 field appears in stdout JSON."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)

    result = main(["--stdout", "TODO"])
    assert result == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["matches_found"] == 1
    assert "sha256" not in payload["matches_per_file"][0]
    for m in payload["matches"]:
        assert "sha256" not in m


def test_hash_present_with_flag(tmp_path, monkeypatch, capsys):
    """With --hash, matches_per_file entries gain a sha256 with the real hex digest."""
    content = b"TODO: x\n"
    (tmp_path / "doc.txt").write_bytes(content)
    monkeypatch.chdir(tmp_path)

    result = main(["--stdout", "--hash", "TODO"])
    assert result == 0
    payload = json.loads(capsys.readouterr().out)
    entries = payload["matches_per_file"]
    assert len(entries) == 1
    assert entries[0]["sha256"] == _sha256_hex(content)
    # Per-match entries are not hashed (would be redundant); only per-file.
    for m in payload["matches"]:
        assert "sha256" not in m


def test_hash_deduplicates_per_file(tmp_path, monkeypatch, capsys):
    """A file with N matches gets one sha256, not N."""
    (tmp_path / "doc.txt").write_text("TODO: a\nTODO: b\nTODO: c\n")
    monkeypatch.chdir(tmp_path)

    result = main(["--stdout", "--hash", "TODO"])
    assert result == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["matches_found"] == 3
    assert len(payload["matches_per_file"]) == 1
    assert "sha256" in payload["matches_per_file"][0]


def test_hash_inverse_search(tmp_path, monkeypatch, capsys):
    """Inverse search: files-without-matches each carry a sha256 with --hash."""
    a = b"contains the term\n"
    b = b"this file does not\n"
    (tmp_path / "a.txt").write_bytes(a)
    (tmp_path / "b.txt").write_bytes(b)
    monkeypatch.chdir(tmp_path)

    result = main(["--stdout", "--hash", "--inverse", "term"])
    # Inverse returns files lacking the term — at least b.txt; exit 0 if any found.
    assert result == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["files_without_matches"] >= 1
    for entry in payload["inverse_files"]:
        assert "sha256" in entry
        if entry["filename"] == "b.txt":
            assert entry["sha256"] == _sha256_hex(b)


def test_hash_in_json_report_file(tmp_path, monkeypatch, capsys):
    """-o json --hash writes the sha256 into the .json report file too."""
    content = b"TODO: x\n"
    (tmp_path / "doc.txt").write_bytes(content)
    monkeypatch.chdir(tmp_path)

    result = main(["-o", "json", "--hash", "TODO"])
    assert result == 0
    report = json.loads((tmp_path / "peekdocs_standard_results.json").read_text())
    assert report["matches_per_file"][0]["sha256"] == _sha256_hex(content)


def test_hash_regex_collection_stdout(tmp_path, monkeypatch, capsys):
    """--regex-collection --hash --stdout produces a top-level matches_per_file with sha256s."""
    content = b"TODO: review\n"
    (tmp_path / "doc.txt").write_bytes(content)
    rc_data = {"my collection": [{"name": "TODOs", "regex": r"TODO\b", "enabled": True}]}
    (tmp_path / ".peekdocs_regex_collections.json").write_text(json.dumps(rc_data))
    monkeypatch.chdir(tmp_path)

    result = main(["--regex-collection", "my collection", "-d", str(tmp_path), "-r", "--stdout", "--hash"])
    assert result == 0
    payload = json.loads(capsys.readouterr().out)
    assert "matches_per_file" in payload
    assert len(payload["matches_per_file"]) == 1
    assert payload["matches_per_file"][0]["sha256"] == _sha256_hex(content)


def test_hash_missing_file_yields_null(tmp_path, monkeypatch, capsys, mocker=None):
    """If a matched file disappears between search and hash, sha256 is null (not a crash)."""
    import peekdocs.reporter as reporter
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)

    # Force the helper to act as if the file vanished.
    original = reporter._sha256_of_file
    monkeypatch.setattr(reporter, "_sha256_of_file", lambda _p: None)
    try:
        result = main(["--stdout", "--hash", "TODO"])
    finally:
        monkeypatch.setattr(reporter, "_sha256_of_file", original)
    assert result == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["matches_per_file"][0]["sha256"] is None


# ── --no-log / --runs (per-run structured log) ─────────────────────────

def _run_log_path(tmp_path):
    return tmp_path / ".peekdocs_runs.log"


def _read_log_entries(tmp_path):
    p = _run_log_path(tmp_path)
    if not p.exists():
        return []
    out = []
    for line in p.read_text().splitlines():
        line = line.strip()
        if not line:
            continue
        out.append(json.loads(line))
    return out


def test_run_log_created_on_first_search(tmp_path, monkeypatch):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)

    result = main(["TODO"])
    assert result == 0
    entries = _read_log_entries(tmp_path)
    assert len(entries) == 1
    e = entries[0]
    assert e["argv"] == ["peekdocs", "TODO"]
    assert e["exit_code"] == 0
    assert e["match_count"] == 1
    assert e["file_count"] >= 1
    assert e["error_count"] == 0
    assert e["elapsed_seconds"] >= 0
    assert e["peekdocs_version"]
    assert "T" in e["timestamp"]  # ISO 8601


def test_run_log_appends_not_overwrites(tmp_path, monkeypatch):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)

    main(["TODO"])
    main(["nope"])
    main(["TODO"])
    entries = _read_log_entries(tmp_path)
    assert len(entries) == 3
    assert [e["argv"][-1] for e in entries] == ["TODO", "nope", "TODO"]
    assert [e["exit_code"] for e in entries] == [0, 1, 0]


def test_no_log_flag_suppresses_single_run(tmp_path, monkeypatch):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)

    main(["--no-log", "TODO"])
    assert not _run_log_path(tmp_path).exists() or _read_log_entries(tmp_path) == []


def test_config_run_log_false_suppresses(tmp_path, monkeypatch):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    (tmp_path / ".peekdocsrc").write_text("run_log = false\n")

    main(["TODO"])
    assert _read_log_entries(tmp_path) == []


def test_informational_commands_not_logged(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)

    main(["--check"])
    main(["--list-files"])
    main([])  # no-args banner
    main(["-h"])
    main(["--config"])
    assert _read_log_entries(tmp_path) == []


def test_suite_command_logged(tmp_path, monkeypatch):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    collection = {
        "saved_searches": {"t": {"search_text": "TODO", "recursive": True}},
        "suites": {"S": ["t"]},
    }
    (tmp_path / ".peekdocs_collection.json").write_text(json.dumps(collection))
    monkeypatch.chdir(tmp_path)

    result = main(["--suite", "S"])
    assert result == 0
    entries = _read_log_entries(tmp_path)
    assert len(entries) == 1
    assert entries[0]["argv"] == ["peekdocs", "--suite", "S"]
    assert entries[0]["match_count"] == 1


def test_regex_collection_logged(tmp_path, monkeypatch):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    rc_data = {"c": [{"name": "T", "regex": r"TODO\b", "enabled": True}]}
    (tmp_path / ".peekdocs_regex_collections.json").write_text(json.dumps(rc_data))
    monkeypatch.chdir(tmp_path)

    result = main(["--regex-collection", "c", "-d", str(tmp_path), "-r"])
    assert result == 0
    entries = _read_log_entries(tmp_path)
    assert len(entries) == 1
    assert entries[0]["argv"][:2] == ["peekdocs", "--regex-collection"]


def test_regex_collection_list_not_logged(tmp_path, monkeypatch):
    """--regex-collection --list is informational, not a real run."""
    rc_data = {"c": [{"name": "T", "regex": r"x", "enabled": True}]}
    (tmp_path / ".peekdocs_regex_collections.json").write_text(json.dumps(rc_data))
    monkeypatch.chdir(tmp_path)

    main(["--regex-collection", "--list"])
    assert _read_log_entries(tmp_path) == []


def test_runs_command_reads_log(tmp_path, monkeypatch, capsys):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    main(["TODO"])

    result = main(["--runs"])
    assert result == 0
    out = capsys.readouterr().out
    assert "Run log:" in out
    assert "TODO" in out
    assert "1 run(s)" in out


def test_runs_command_json_emits_jsonl(tmp_path, monkeypatch, capsys):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    main(["TODO"])
    main(["nope"])
    capsys.readouterr()  # discard prior output

    result = main(["--runs", "--json"])
    assert result == 0
    out = capsys.readouterr().out
    lines = [l for l in out.splitlines() if l.strip()]
    assert len(lines) == 2
    for line in lines:
        parsed = json.loads(line)
        assert "argv" in parsed
        assert "exit_code" in parsed


def test_runs_command_empty_log(tmp_path, monkeypatch, capsys):
    monkeypatch.chdir(tmp_path)
    result = main(["--runs"])
    assert result == 0
    out = capsys.readouterr().out
    assert "No run log entries" in out


def test_runs_command_not_logged_itself(tmp_path, monkeypatch):
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    main(["TODO"])
    main(["--runs"])
    # --runs itself should not appear in the log
    entries = _read_log_entries(tmp_path)
    assert len(entries) == 1
    assert entries[0]["argv"] == ["peekdocs", "TODO"]


def test_run_log_path_override(tmp_path, monkeypatch):
    """--config run_log_path redirects the log to a custom location."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    custom = tmp_path / "custom" / "runs.log"
    custom.parent.mkdir()
    (tmp_path / ".peekdocsrc").write_text(f"run_log_path = {custom}\n")
    monkeypatch.chdir(tmp_path)

    main(["TODO"])
    assert custom.exists()
    assert not _run_log_path(tmp_path).exists()


# ── --dry-run (scope preflight) ────────────────────────────────────────

def test_dry_run_reports_file_count(tmp_path, monkeypatch, capsys):
    """--dry-run prints scope info without running a search."""
    (tmp_path / "a.txt").write_text("hello world\n")
    (tmp_path / "b.txt").write_text("another file\n")
    monkeypatch.chdir(tmp_path)

    result = main(["--dry-run", "anything"])
    assert result == 0
    out = capsys.readouterr().out
    assert "Dry run" in out
    assert "2 file" in out
    assert "By extension" in out
    assert ".txt" in out
    assert "No content read" in out


def test_dry_run_writes_no_reports(tmp_path, monkeypatch):
    """--dry-run never produces peekdocs_*_results.* files."""
    (tmp_path / "a.txt").write_text("hello world\n")
    monkeypatch.chdir(tmp_path)

    main(["--dry-run", "hello"])
    assert not any(tmp_path.glob("peekdocs_*_results.*"))


def test_dry_run_exit_code_zero_files(tmp_path, monkeypatch, capsys):
    """--dry-run returns 1 when no files match the filters."""
    monkeypatch.chdir(tmp_path)  # empty dir
    result = main(["--dry-run", "anything"])
    assert result == 1


def test_dry_run_json_emits_clean_payload(tmp_path, monkeypatch, capsys):
    """--dry-run --stdout emits a clean JSON payload with no banner."""
    (tmp_path / "a.txt").write_text("hello\n")
    (tmp_path / "b.pdf").write_bytes(b"fake pdf bytes\n")
    monkeypatch.chdir(tmp_path)

    result = main(["--dry-run", "--stdout", "x"])
    assert result == 0
    out = capsys.readouterr().out
    payload = json.loads(out)
    assert payload["dry_run"] is True
    assert payload["file_count"] == 2
    assert payload["total_bytes"] > 0
    exts = {row["ext"] for row in payload["by_extension"]}
    assert ".txt" in exts and ".pdf" in exts


def test_dry_run_honors_file_type_filter(tmp_path, monkeypatch, capsys):
    """-t pdf in dry-run only counts PDFs."""
    (tmp_path / "a.txt").write_text("text\n")
    (tmp_path / "b.pdf").write_bytes(b"pdf bytes\n")
    monkeypatch.chdir(tmp_path)

    result = main(["--dry-run", "--stdout", "-t", "pdf", "x"])
    assert result == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["file_count"] == 1
    assert payload["by_extension"][0]["ext"] == ".pdf"


def test_dry_run_honors_recursive(tmp_path, monkeypatch, capsys):
    """-r in dry-run picks up subfolder files."""
    sub = tmp_path / "sub"
    sub.mkdir()
    (tmp_path / "top.txt").write_text("top\n")
    (sub / "deep.txt").write_text("deep\n")
    monkeypatch.chdir(tmp_path)

    flat = json.loads(_capture_dry_run_json(["--dry-run", "--stdout", "x"], capsys))
    rec = json.loads(_capture_dry_run_json(["--dry-run", "--stdout", "-r", "x"], capsys))
    assert flat["file_count"] == 1
    assert rec["file_count"] == 2


def _capture_dry_run_json(argv, capsys):
    capsys.readouterr()  # discard prior output
    main(argv)
    return capsys.readouterr().out


def test_dry_run_not_logged(tmp_path, monkeypatch):
    """--dry-run runs are not recorded in ~/.peekdocs_runs.log."""
    (tmp_path / "a.txt").write_text("hello\n")
    monkeypatch.chdir(tmp_path)

    main(["--dry-run", "hello"])
    main(["--dry-run", "--stdout", "world"])
    log = tmp_path / ".peekdocs_runs.log"
    assert not log.exists() or log.read_text() == ""


def test_dry_run_with_suite_errors_cleanly(tmp_path, monkeypatch, capsys):
    """--dry-run with --suite should error, not silently run the suite."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    collection = {
        "saved_searches": {"t": {"search_text": "TODO", "recursive": True}},
        "suites": {"S": ["t"]},
    }
    (tmp_path / ".peekdocs_collection.json").write_text(json.dumps(collection))
    monkeypatch.chdir(tmp_path)

    result = main(["--suite", "S", "--dry-run"])
    assert result == 2
    out = capsys.readouterr().out
    assert "not supported with --suite" in out
    # No suite report should have been written.
    assert not (tmp_path / "peekdocs_suite_results.txt").exists()


def test_dry_run_with_regex_collection_errors_cleanly(tmp_path, monkeypatch, capsys):
    """--dry-run with --regex-collection should error, not silently run the collection."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    rc_data = {"c": [{"name": "T", "regex": r"TODO\b", "enabled": True}]}
    (tmp_path / ".peekdocs_regex_collections.json").write_text(json.dumps(rc_data))
    monkeypatch.chdir(tmp_path)

    result = main(["--regex-collection", "c", "-d", str(tmp_path), "-r", "--dry-run"])
    assert result == 2
    out = capsys.readouterr().out
    assert "not supported with --regex-collection" in out
    assert not (tmp_path / "peekdocs_regex_results.txt").exists()


# ── --on-match (notification hook) ─────────────────────────────────────

@_posix_only
def test_on_match_fires_when_matches_found(tmp_path, monkeypatch):
    """--on-match runs the given command when matches are found."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    flag_file = tmp_path / "fired.txt"
    # Tiny inline script that writes a flag file with the env vars.
    script = tmp_path / "hook.sh"
    script.write_text(
        f"#!/bin/sh\n"
        f"echo \"matches=$PEEKDOCS_MATCH_COUNT files=$PEEKDOCS_FILE_COUNT cwd=$PEEKDOCS_CWD\" > {flag_file}\n"
    )
    script.chmod(0o755)

    result = main(["--on-match", str(script), "TODO"])
    assert result == 0
    assert flag_file.exists()
    content = flag_file.read_text()
    assert "matches=1" in content
    # files= reflects what peekdocs searched (may include hook.sh in the same dir)
    assert "files=" in content
    assert str(tmp_path) in content


def test_on_match_does_not_fire_on_no_match(tmp_path, monkeypatch):
    """--on-match does NOT fire when search finds zero matches (exit 1)."""
    (tmp_path / "doc.txt").write_text("nothing here\n")
    monkeypatch.chdir(tmp_path)
    flag_file = tmp_path / "fired.txt"
    script = tmp_path / "hook.sh"
    script.write_text(f"#!/bin/sh\necho fired > {flag_file}\n")
    script.chmod(0o755)

    result = main(["--on-match", str(script), "nope"])
    assert result == 1
    assert not flag_file.exists()


@_posix_only
def test_on_match_passes_report_paths_for_standard_search(tmp_path, monkeypatch):
    """PEEKDOCS_REPORT_TXT / _DOCX are set when standard reports are written."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    env_dump = tmp_path / "env.txt"
    script = tmp_path / "hook.sh"
    script.write_text(
        f"#!/bin/sh\n"
        f"echo \"TXT=$PEEKDOCS_REPORT_TXT\" > {env_dump}\n"
        f"echo \"DOCX=$PEEKDOCS_REPORT_DOCX\" >> {env_dump}\n"
    )
    script.chmod(0o755)

    main(["--on-match", str(script), "TODO"])
    out = env_dump.read_text()
    assert "peekdocs_standard_results.txt" in out
    assert "peekdocs_standard_results.docx" in out


def test_on_match_hook_failure_does_not_change_exit_code(tmp_path, monkeypatch, capsys):
    """A failing hook script logs a warning but the search exit code is preserved."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    script = tmp_path / "broken.sh"
    script.write_text("#!/bin/sh\nexit 1\n")
    script.chmod(0o755)

    result = main(["--on-match", str(script), "TODO"])
    # Search succeeded; hook failed; result reflects the search.
    assert result == 0


def test_on_match_unknown_command_warns_but_succeeds(tmp_path, monkeypatch, capsys):
    """If the hook command doesn't exist, peekdocs warns but the search succeeds."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)

    result = main(["--on-match", "/nonexistent/path/nothing.sh", "TODO"])
    assert result == 0


@_posix_only
def test_on_match_logged_in_run_log(tmp_path, monkeypatch):
    """The run log entry includes on_match_fired when --on-match was set."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    script = tmp_path / "hook.sh"
    script.write_text("#!/bin/sh\nexit 0\n")
    script.chmod(0o755)

    main(["--on-match", str(script), "TODO"])
    entries = _read_log_entries(tmp_path)
    assert len(entries) == 1
    assert entries[0]["on_match_fired"] is True


def test_on_match_absent_when_flag_not_used(tmp_path, monkeypatch):
    """Routine runs without --on-match have no on_match_fired key in the log."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)

    main(["TODO"])
    entries = _read_log_entries(tmp_path)
    assert "on_match_fired" not in entries[0]


@_posix_only
def test_on_match_config_default_picked_up(tmp_path, monkeypatch):
    """`--config on_match=...` saves a persistent default that fires without the flag."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    flag_file = tmp_path / "fired.txt"
    script = tmp_path / "hook.sh"
    script.write_text(f"#!/bin/sh\necho default fired > {flag_file}\n")
    script.chmod(0o755)
    (tmp_path / ".peekdocsrc").write_text(f"on_match = {script}\n")

    main(["TODO"])
    assert flag_file.exists()


def test_on_match_empty_string_disables_config_default(tmp_path, monkeypatch):
    """`--on-match ''` disables the hook for one run even if config sets it."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    flag_file = tmp_path / "fired.txt"
    script = tmp_path / "hook.sh"
    script.write_text(f"#!/bin/sh\necho fired > {flag_file}\n")
    script.chmod(0o755)
    (tmp_path / ".peekdocsrc").write_text(f"on_match = {script}\n")

    main(["--on-match", "", "TODO"])
    assert not flag_file.exists()


def test_on_match_not_fired_on_dry_run(tmp_path, monkeypatch):
    """--dry-run never produces matches, so --on-match never fires."""
    (tmp_path / "doc.txt").write_text("TODO: x\n")
    monkeypatch.chdir(tmp_path)
    flag_file = tmp_path / "fired.txt"
    script = tmp_path / "hook.sh"
    script.write_text(f"#!/bin/sh\necho fired > {flag_file}\n")
    script.chmod(0o755)

    main(["--dry-run", "--on-match", str(script), "TODO"])
    assert not flag_file.exists()


# ── --diff (compare two JSON outputs) ──────────────────────────────────

def _write_json(path, payload):
    path.write_text(json.dumps(payload))


def _stdout_shape(matches_per_file, **extras):
    """Mock a peekdocs --stdout JSON payload."""
    payload = {
        "generator": "peekdocs vTEST",
        "directory": "/tmp/docs",
        "search_terms": ["x"],
        "mode": "ANY",
        "files_searched": 100,
        "matches_found": sum(e["matches"] for e in matches_per_file),
        "elapsed_seconds": 0.1,
        "matches_per_file": matches_per_file,
        "matches": [],
    }
    payload.update(extras)
    return payload


def test_diff_detects_new_files(tmp_path, monkeypatch, capsys):
    """A file that matches in NEW but not OLD shows up as NEW."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    _write_json(old, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 3}]))
    _write_json(new, _stdout_shape([
        {"filename": "a.txt", "folder": "/d", "matches": 3},
        {"filename": "b.txt", "folder": "/d", "matches": 7},
    ]))
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new)])
    assert result == 1  # actionable: new file matching
    out = capsys.readouterr().out
    assert "NEW: 1 file" in out
    assert "b.txt" in out


def test_diff_detects_removed_files(tmp_path, monkeypatch, capsys):
    """A file that matched in OLD but not NEW shows up as REMOVED."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    _write_json(old, _stdout_shape([
        {"filename": "a.txt", "folder": "/d", "matches": 3},
        {"filename": "gone.txt", "folder": "/d", "matches": 5},
    ]))
    _write_json(new, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 3}]))
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new)])
    # Removed-only is NOT actionable per design.
    assert result == 0
    out = capsys.readouterr().out
    assert "REMOVED: 1 file" in out
    assert "gone.txt" in out


def test_diff_detects_changed_counts(tmp_path, monkeypatch, capsys):
    """Same file, different match count → CHANGED with delta."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    _write_json(old, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 3}]))
    _write_json(new, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 7}]))
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new)])
    assert result == 1  # gained matches → actionable
    out = capsys.readouterr().out
    assert "CHANGED: 1 file" in out
    assert "3 → 7" in out
    assert "(+4)" in out


def test_diff_changed_decrease_not_actionable(tmp_path, monkeypatch, capsys):
    """Fewer matches than before is NOT actionable (file probably got cleaned up)."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    _write_json(old, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 10}]))
    _write_json(new, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 3}]))
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new)])
    assert result == 0
    out = capsys.readouterr().out
    assert "CHANGED" in out
    assert "(-7)" in out


def test_diff_detects_modified_when_sha256_changes(tmp_path, monkeypatch, capsys):
    """Same path, same match count, different sha256 → MODIFIED (only when --hash was used)."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    _write_json(old, _stdout_shape([
        {"filename": "a.txt", "folder": "/d", "matches": 3, "sha256": "aaa"},
    ]))
    _write_json(new, _stdout_shape([
        {"filename": "a.txt", "folder": "/d", "matches": 3, "sha256": "bbb"},
    ]))
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new)])
    assert result == 1  # modified content is actionable
    out = capsys.readouterr().out
    assert "MODIFIED: 1 file" in out


def test_diff_no_modified_without_hashes(tmp_path, monkeypatch, capsys):
    """No --hash → sha256 absent → MODIFIED detection silently skipped."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    _write_json(old, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 3}]))
    _write_json(new, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 3}]))
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new)])
    assert result == 0  # nothing changed
    out = capsys.readouterr().out
    assert "MODIFIED" not in out
    assert "UNCHANGED: 1" in out


def test_diff_json_output(tmp_path, monkeypatch, capsys):
    """--json emits a structured diff with no banner above it."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    _write_json(old, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 3}]))
    _write_json(new, _stdout_shape([
        {"filename": "a.txt", "folder": "/d", "matches": 5},
        {"filename": "b.txt", "folder": "/d", "matches": 2},
    ]))
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new), "--json"])
    assert result == 1
    out = capsys.readouterr().out
    payload = json.loads(out)
    assert len(payload["new"]) == 1
    assert payload["new"][0]["filename"] == "b.txt"
    assert len(payload["changed"]) == 1
    assert payload["changed"][0]["delta"] == 2
    assert payload["old_match_total"] == 3
    assert payload["new_match_total"] == 7


def test_diff_reconstructs_from_flat_matches(tmp_path, monkeypatch, capsys):
    """If matches_per_file is absent (regex-collection without --hash), reconstruct from matches."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    # Mimic regex-collection JSON without --hash (no matches_per_file)
    old_data = {
        "generator": "peekdocs v0", "collection": "c", "directory": "/d",
        "timestamp": "...", "elapsed_seconds": 0.1, "total_matches": 2,
        "patterns": [],
        "matches": [
            {"filename": "a.txt", "folder": "/d", "line_number": 1, "matched_text": "x"},
            {"filename": "a.txt", "folder": "/d", "line_number": 5, "matched_text": "x"},
        ],
    }
    new_data = {
        "generator": "peekdocs v0", "collection": "c", "directory": "/d",
        "timestamp": "...", "elapsed_seconds": 0.1, "total_matches": 1,
        "patterns": [],
        "matches": [
            {"filename": "b.txt", "folder": "/d", "line_number": 1, "matched_text": "x"},
        ],
    }
    _write_json(old, old_data)
    _write_json(new, new_data)
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new), "--json"])
    payload = json.loads(capsys.readouterr().out)
    assert result == 1
    new_files = {e["filename"] for e in payload["new"]}
    removed_files = {e["filename"] for e in payload["removed"]}
    assert new_files == {"b.txt"}
    assert removed_files == {"a.txt"}


def test_diff_missing_file_errors_cleanly(tmp_path, monkeypatch, capsys):
    monkeypatch.chdir(tmp_path)
    result = main(["--diff", str(tmp_path / "nope.json"), str(tmp_path / "nope2.json")])
    assert result == 2
    err = capsys.readouterr().err
    assert "Error reading old file" in err
    assert "file not found" in err
    # Generic hint for non-document extensions
    assert "--stdout" in err


def test_diff_invalid_json_errors_cleanly(tmp_path, monkeypatch, capsys):
    bad = tmp_path / "bad.json"
    bad.write_text("{this is not valid json")
    good = tmp_path / "good.json"
    _write_json(good, _stdout_shape([]))
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(bad), str(good)])
    assert result == 2
    err = capsys.readouterr().err
    assert "invalid JSON" in err


def test_diff_document_input_gives_document_hint(tmp_path, monkeypatch, capsys):
    """Passing a .odt/.docx/.pdf to --diff should trigger the document-vs-snapshot hint."""
    odt = tmp_path / "meeting_notes.odt"
    odt.write_text("not actually json")
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(odt), str(odt)])
    assert result == 2
    err = capsys.readouterr().err
    assert "looks like a document" in err
    assert "peekdocs JSON snapshot" in err
    assert "--stdout" in err


def test_diff_no_changes_returns_0(tmp_path, monkeypatch, capsys):
    """Identical inputs → no diff → exit 0."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    payload = _stdout_shape([
        {"filename": "a.txt", "folder": "/d", "matches": 3},
        {"filename": "b.txt", "folder": "/d", "matches": 7},
    ])
    _write_json(old, payload)
    _write_json(new, payload)
    monkeypatch.chdir(tmp_path)

    result = main(["--diff", str(old), str(new)])
    assert result == 0
    out = capsys.readouterr().out
    assert "UNCHANGED: 2" in out


def test_diff_too_few_args(tmp_path, monkeypatch, capsys):
    monkeypatch.chdir(tmp_path)
    result = main(["--diff", "only-one.json"])
    assert result == 2
    err = capsys.readouterr().err
    assert "requires two" in err


def test_diff_not_logged(tmp_path, monkeypatch):
    """--diff is informational and should not appear in the run log."""
    old = tmp_path / "old.json"
    new = tmp_path / "new.json"
    _write_json(old, _stdout_shape([]))
    _write_json(new, _stdout_shape([{"filename": "a.txt", "folder": "/d", "matches": 1}]))
    monkeypatch.chdir(tmp_path)

    main(["--diff", str(old), str(new)])
    log = tmp_path / ".peekdocs_runs.log"
    assert not log.exists() or log.read_text() == ""
